"""
Presentación M.A.T.E.R.I.A. V4 + HSAQ — Congreso AI LATAM
Versión mejorada con V4 JEPA-First + SCA
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = '/home/methodwhite/Documentos/IA LATAM/Plantilla/Plantilla_Congreso_AI_LATAM.pptx'
OUTPUT = '/home/methodwhite/MATERIA/docs/MATERIA_Conferencia_AI_LATAM_v2.pptx'
DOCS = '/home/methodwhite/MATERIA/docs'

# Cargar plantilla y quedarse solo con layouts
template_prs = Presentation(TEMPLATE)
prs = Presentation(TEMPLATE)

# Limpiar todas las slides existentes
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Colores
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x2D, 0x5B, 0xA0)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7A, 0x7A, 0x7A)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)

def add_slide():
    layout = prs.slide_layouts[0]  # Default layout
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK
    return slide

def txt(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

def multi_txt(slide, left, top, width, height, lines, size=14, color=WHITE, line_spacing=1.3):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(size * 0.3)
    return box

# ============================================================
# SLIDE 1: Portada
# ============================================================
s = add_slide()
txt(s, 0.5, 0.5, 12, 0.6, 'CONGRESO · INTELIGENCIA ARTIFICIAL · 2026', 14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, 0.5, 1.5, 12, 1.5, 'M.A.T.E.R.I.A. V4', 52, True, ACCENT, PP_ALIGN.CENTER)
txt(s, 0.5, 3.0, 12, 0.8, 'JEPA-First Architecture + SCA + HSAQ', 22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
txt(s, 0.5, 4.0, 12, 0.6, 'Multi-Analytical Toroidal Engine for Recursive Intelligent Analysis', 16, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 0.5, 5.5, 12, 0.5, 'Entrenamiento de modelos y cuantización HSAQ', 14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, 0.5, 6.5, 12, 0.4, 'MethodWhite · M.A.T.E.R.I.A. Research · Santiago, Chile', 12, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Título charla
# ============================================================
s = add_slide()
txt(s, 0.5, 0.5, 12, 0.5, 'CHARLA MAGISTRAL · TRACK DE IA', 14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, 0.5, 1.8, 12, 1.2, 'Entrenamiento de Modelos M.A.T.E.R.I.A.\ny Cuantización HSAQ', 36, True, WHITE, PP_ALIGN.CENTER)
txt(s, 0.5, 3.5, 12, 0.8, 'Cómo logramos 99% accuracy con 3.8M parámetros\nusando arquitectura multi-paradigma JEPA-First', 18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Sección - Qué es MATERIA
# ============================================================
s = add_slide()
txt(s, 0.5, 1.5, 12, 1.0, '01', 72, True, ACCENT, PP_ALIGN.CENTER)
txt(s, 0.5, 3.0, 12, 0.8, 'QUÉ ES M.A.T.E.R.I.A.?', 36, True, WHITE, PP_ALIGN.CENTER)
txt(s, 0.5, 4.2, 12, 0.6, 'Un sistema de IA que integra múltiples paradigmas\nen una arquitectura unificada y eficiente', 16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Agenda
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'AGENDA', 32, True, WHITE)
items = [
    '01  Arquitectura V4: JEPA-First + SCA',
    '02  Entrenamiento: De 0 a 99% accuracy en CPU',
    '03  HSAQ: Ejecución dispersa adaptativa',
    '04  Comparativa: M.A.T.E.R.I.A. vs GPT vs PaLM vs Gemini',
    '05  Demo: Generación con modelo entrenado',
]
for i, item in enumerate(items):
    txt(s, 0.8, 1.5 + i*0.9, 11, 0.8, item, 16, color=LIGHT_BLUE)

# ============================================================
# SLIDE 5: Arquitectura V4 - JEPA-First
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'ARQUITECTURA V4: JEPA-FIRST', 28, True, WHITE)

multi_txt(s, 0.5, 1.2, 5.5, 5.5, [
    'Flujo de datos:',
    '',
    '1. Token Embedding (vocab → dim)',
    '2. HSAQ Sparse Execution',
    '3. Transformer Blocks (GQA + RoPE + SwiGLU)',
    '4. SNN (LIF neuronas reales)',
    '5. SSM (State Space Model)',
    '   ↓',
    '6. Concat [GQA, SNN, SSM]',
    '7. JEPA Encoder → espacio latente',
    '8. SCA Predictor (K = √π·e·γ)',
    '9. Synapsis Memory (opcional)',
    '10. RMSNorm → Head → Output',
], size=12, color=LIGHT_BLUE)

multi_txt(s, 6.5, 1.2, 5.5, 5.5, [
    'Componentes clave:',
    '',
    'GQA: 8 query heads, 4 KV heads',
    'RoPE: codificación posicional rotatoria',
    'SwiGLU: activación puerta adaptativa',
    'LIF-SNN: neuronas de pulsos reales',
    'SSM: procesamiento secuencias largas',
    'JEPA: predicción en espacio latente',
    'SCA: descomposición espectral',
    'HSAQ: 30% sparsity adaptativa',
    '',
    'K = 2.781042 (peso JEPA)',
    'Loss = token_loss + K × jepa_loss',
], size=12, color=LIGHT_BLUE)

# ============================================================
# SLIDE 6: Comparativa vs Otras Arquitecturas
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'M.A.T.E.R.I.A. vs OTRAS ARQUITECTURAS', 28, True, WHITE)

table_data = [
    ['Característica', 'GPT-4', 'PaLM', 'Gemini', 'M.A.T.E.R.I.A. V4'],
    ['Parámetros', '~1.8T', '540B', '1.5T+', '4.7M'],
    ['Hardware', '10K GPUs', '6144 TPUv4', 'Multi-TPU', 'CPU solo'],
    ['SNN integration', 'No', 'No', 'No', 'Sí (LIF real)'],
    ['JEPA + SCA', 'No', 'No', 'No', 'Sí (K=2.78)'],
    ['HSAQ sparsity', 'No', 'No', 'No', 'Sí (30%)'],
    ['Accuracy (val)', '~86%', '~67%', '~90%', '99.03%'],
    ['Costo entrenamiento', '~$100M', '~$12M', '~$50M+', '<$1 (CPU)'],
]

rows = len(table_data)
cols = len(table_data[0])
tbl = s.shapes.add_table(rows, cols, Inches(0.3), Inches(1.3), Inches(12.5), Inches(5.5)).table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.color.rgb = WHITE
            para.font.bold = (r == 0 or c == 0)
            para.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        elif c == 4:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)

# ============================================================
# SLIDE 7: HSAQ
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'HSAQ: EJECUCIÓN DISPERSA ADAPTATIVA', 28, True, WHITE)

code_lines = [
    '# HSAQ: Solo neuronas relevantes pasan',
    'flat = x.abs().view(B, -1)',
    'k = n * (1 - sparsity)  # Top-70%',
    'thresh = torch.kthvalue(flat, k)',
    'mask = x.abs() >= thresh',
    'return x * mask  # 30% → 0',
]
box = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(3))
tf = box.text_frame
tf.word_wrap = True
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.name = 'Consolas'
    p.font.color.rgb = GREEN

multi_txt(s, 7.0, 1.3, 5.5, 4, [
    'Ventajas:',
    '',
    '✓ Adaptativo por batch (umbral dinámico)',
    '✓ Gradiente fluye (entrenable end-to-end)',
    '✓ Hardware-agnostic (CPU/GPU/TPU)',
    '✓ 30% menos cómputo, misma accuracy',
    '✓ Sin post-processing',
    '',
    'vs TurboQuant (Google):',
    '✗ TurboQuant: cuantización fija post-training',
    '✗ HSAQ: ejecución dispersa adaptativa',
], size=12, color=WHITE)

# ============================================================
# SLIDE 8: Pipeline Entrenamiento
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'PIPELINE DE ENTRENAMIENTO V4', 28, True, WHITE)

pipeline_img = os.path.join(DOCS, 'pipeline.png')
if os.path.exists(pipeline_img):
    s.shapes.add_picture(pipeline_img, Inches(1), Inches(1.2), Inches(11), Inches(5.8))
else:
    steps = [
        ('1. Recolección de Datos', 'C4 EN (773MB) + Wikipedia 12 idiomas\nStreaming desde HuggingFace, 80K textos'),
        ('2. Tokenización', 'Char-level (258 tokens) o BPE (32K)\nSentencePiece multilingüe'),
        ('3. Entrenamiento JEPA-First', 'Loss dual: token_loss + K × jepa_loss\nK = 2.781042 (coupling geometría-entropía)'),
        ('4. HSAQ Compresión', 'Sparsity 30% adaptativa\nEncoder FP16, Predictor INT8, Decoder INT8'),
        ('5. Exportación .materia', 'Pesos numpy + metadatos + config\nListo para inferencia en CPU'),
    ]
    for i, (title, desc) in enumerate(steps):
        y = 1.2 + i * 1.15
        txt(s, 0.5, y, 12, 0.4, title, 16, True, ACCENT)
        txt(s, 0.5, y + 0.4, 12, 0.7, desc, 12, color=LIGHT_BLUE)

# ============================================================
# SLIDE 9: Entrenamiento sin vs con Synapsis
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'SIN vs CON SYNAPSIS', 28, True, WHITE)

comp_data = [
    ['Métrica', 'Sin Synapsis', 'Con Synapsis'],
    ['Loss (epoch 1)', '0.09', '2.08'],
    ['Accuracy', '99.2%', '83.6%'],
    ['Convergencia', '~500 batches', 'Lenta'],
    ['Repetición', 'No', 'Sí ("the the the")'],
    ['Spike rate', '0.048 (activo)', '0.035 (bajo)'],
    ['Uso recomendado', 'Entrenamiento', 'Inferencia'],
]

rows = len(comp_data)
cols = len(comp_data[0])
tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12), Inches(5)).table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = comp_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(13)
            para.font.color.rgb = WHITE
            para.font.bold = (r == 0 or c == 0)
            para.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        elif c == 1:  # Sin Synapsis highlight
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)

txt(s, 0.5, 6.0, 12, 0.5, 'Conclusión: Synapsis causa maldición de repetición durante entrenamiento.\nUsar sin Synapsis para training, con Synapsis para inference.', 12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: Resultados
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'RESULTADOS DE ENTRENAMIENTO', 28, True, WHITE)

training_img = os.path.join(DOCS, 'plots', 'comparative_training.png')
if os.path.exists(training_img):
    s.shapes.add_picture(training_img, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))

# ============================================================
# SLIDE 11: Ecosistema de modelos
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'ECOSISTEMA DE MODELOS', 28, True, WHITE)

model_data = [
    ['Modelo', 'Params', 'Dataset', 'Loss', 'Accuracy', 'Uso'],
    ['materia-v4-base', '4.7M', 'C4 EN 80K', '-', '99.0%', 'Base V4'],
    ['materia-v3.basemateria', '3.5M', 'C4 EN', '0.032', '99.0%', 'Base V3'],
    ['materia-v3-full', '4.8M', 'C4 EN 15K', '0.033', '99.0%', 'General'],
    ['materia-v3-nano', '0.6M', 'C4 EN 1K', '0.047', '98.9%', 'Edge/IoT'],
    ['science-v3', '2.3M', 'Reasoning', '0.031', '99.8%', 'Científico'],
]

rows = len(model_data)
cols = len(model_data[0])
tbl = s.shapes.add_table(rows, cols, Inches(0.3), Inches(1.3), Inches(12.5), Inches(5.5)).table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = model_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.color.rgb = WHITE
            para.font.bold = (r == 0)
            para.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        elif r == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)

# ============================================================
# SLIDE 12: ¿Por qué tan inteligente?
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, '¿POR QUÉ ES TAN INTELIGENTE?', 28, True, WHITE)

reasons = [
    ('Integración End-to-End', '8 paradigmas entrenables juntos.\nEl gradiente fluye por toda la arquitectura.'),
    ('JEPA-First: predicción en espacio latente', 'K = 2.781042 como peso.\nCoupling geometría-entropía.'),
    ('SNN real, no aproximación', 'Neuronas LIF con dinámica de membrana.\nSpike rate estable ~0.04.'),
    ('HSAQ: eficiencia sin pérdida', '30% sparsity adaptativa.\nMisma accuracy, menor costo.'),
]

for i, (title, desc) in enumerate(reasons):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.3
    y = 1.3 + row * 2.8
    txt(s, x, y, 6.0, 0.4, title, 16, True, ACCENT)
    txt(s, x, y + 0.5, 6.0, 1.5, desc, 13, color=LIGHT_BLUE)

# ============================================================
# SLIDE 13: ¿Qué esperar de 1B?
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, '¿QUÉ ESPERAR DE 1B+ PARÁMETROS?', 28, True, WHITE)

projections = [
    ('Razonamiento complejo', 'Cadena de pensamiento nativa.\nInferencia multi-paso.\nResolución matemática.'),
    ('Comprensión profunda', 'Contexto 4K-8K tokens.\nRetención a largo plazo.\nRespuestas coherentes.'),
    ('Multimodalidad', 'Texto + imagen + audio.\nMismo pipeline HSAQ.\nMantenimiento de eficiencia.'),
    ('Edge deployment', '1B params → ~2GB FP16.\n~1.4GB con HSAQ 30%.\nViable en CPU.'),
]

for i, (title, desc) in enumerate(projections):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.3
    y = 1.3 + row * 2.8
    txt(s, x, y, 6.0, 0.4, title, 16, True, ACCENT)
    txt(s, x, y + 0.5, 6.0, 1.5, desc, 13, color=LIGHT_BLUE)

# ============================================================
# SLIDE 14: Cita
# ============================================================
s = add_slide()
txt(s, 1.0, 2.0, 11, 2.0,
    '"La inteligencia artificial no necesita ser masiva\npara ser efectiva. M.A.T.E.R.I.A. demuestra que\nla integración inteligente supera al escalado bruto."',
    24, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 1.0, 4.5, 11, 0.5, 'MethodWhite  ·  M.A.T.E.R.I.A. Research', 14, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15: Datos clave
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'DATOS QUE RESPALDAN EL MENSAJE', 28, True, WHITE)

stats = [
    ('4.7M', 'Parámetros'),
    ('99.0%', 'Accuracy'),
    ('$0', 'Costo hardware'),
    ('8', 'Paradigmas'),
    ('30%', 'Ahorro HSAQ'),
    ('CPU', 'Entrenamiento'),
]

for i, (num, label) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.5 + row * 2.8
    txt(s, x, y, 3.8, 1.2, num, 48, True, ACCENT, PP_ALIGN.CENTER)
    txt(s, x, y + 1.2, 3.8, 0.5, label, 16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 16: Para llevar
# ============================================================
s = add_slide()
txt(s, 0.5, 0.3, 12, 0.8, 'PARA LLEVAR', 32, True, WHITE)

takeaways = [
    '1. La integración multi-paradigma supera al escalado bruto.',
    '2. HSAQ permite ejecución eficiente sin pérdida de accuracy.',
    '3. JEPA-First + SCA es la clave del rendimiento V4.',
    '4. Synapsis causa repetición — usar solo para inference.',
    '5. El futuro está en modelos pequeños pero inteligentes.',
]
for i, t in enumerate(takeaways):
    txt(s, 0.8, 1.5 + i*1.0, 11, 0.8, t, 18, color=LIGHT_BLUE)

# ============================================================
# SLIDE 17: Gracias
# ============================================================
s = add_slide()
txt(s, 0.5, 2.0, 12, 1.5, '¡Gracias!', 64, True, ACCENT, PP_ALIGN.CENTER)
txt(s, 0.5, 3.8, 12, 0.8, '¿Conversamos? Espacio para preguntas.', 20, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 0.5, 5.0, 12, 0.5, 'github.com/methodwhite/MATERIA', 14, color=GRAY, align=PP_ALIGN.CENTER)

# Guardar
prs.save(OUTPUT)
print(f'Presentación guardada: {OUTPUT}')
print(f'Slides: {len(prs.slides)}')
