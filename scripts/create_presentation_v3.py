#!/usr/bin/env python3
"""
MATERIA V4 + HSAQ — Presentation for Congreso AI LATAM
Built from the official template's color palette and proportions.
Uses only python-pptx shapes (no external images).
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import copy

# ── Paths ──────────────────────────────────────────────────────────────
TEMPLATE_PATH = '/home/methodwhite/Documentos/IA LATAM/Plantilla/Plantilla_Congreso_AI_LATAM.pptx'
OUTPUT_PATH = '/home/methodwhite/MATERIA/outputs/MATERIA_V4_HSAQ_Congreso.pptx'

# ── Template colour palette (from slide 26) ───────────────────────────
NAVY       = RGBColor(0x0B, 0x1A, 0x47)
NAVY_DEEP  = RGBColor(0x07, 0x0F, 0x2E)
NAVY_SURF  = RGBColor(0x16, 0x26, 0x5E)
CIAN       = RGBColor(0x47, 0xB6, 0xFF)
CIAN_SOFT  = RGBColor(0x8F, 0xD3, 0xFF)
CIAN_LIGHT = RGBColor(0x5F, 0xBB, 0xFF)
HIELO      = RGBColor(0xC9, 0xD8, 0xF5)
CLARO      = RGBColor(0xF4, 0xF7, 0xFC)
TINTA      = RGBColor(0x5C, 0x6B, 0x8A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Gold / amber accents (requested by speaker)
AMBER      = RGBColor(0xFF, 0xB3, 0x00)
AMBER_DARK = RGBColor(0xCC, 0x88, 0x00)
AMBER_SOFT = RGBColor(0xFF, 0xD5, 0x66)

# ── Open template (to preserve dimensions) ────────────────────────────
prs = Presentation(TEMPLATE_PATH)
SW = prs.slide_width   # 12191695 EMU ≈ 13.33"
SH = prs.slide_height   # 6858000 EMU ≈ 7.50"

# ── Helpers ───────────────────────────────────────────────────────────

def _emu(inches):
    """Convert inches to EMU."""
    return int(inches * 914400)

def _add_bg(slide, color=NAVY_DEEP):
    """Set slide background to a solid fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def _rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def _circle(slide, left, top, diameter, fill_color, line_color=None):
    """Add a circle (oval)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def _txt(slide, left, top, width, height, text,
         font_name='Calibri', font_size=Pt(16), bold=False,
         color=HIELO, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         word_wrap=True):
    """Add a text box with a single paragraph."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.paragraphs[0].alignment = alignment
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    # Set anchor
    try:
        txbox.text_frame._txBody.bodyPr.set('anchor', {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }[anchor])
    except Exception:
        pass
    return txbox

def _multi_txt(slide, left, top, width, height, lines,
               font_name='Calibri', font_size=Pt(16),
               color=HIELO, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs (list of strings)."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color
        p.space_after = Pt(4)
    return txbox

def _card(slide, left, top, width, height, title, body_lines,
          title_color=CIAN, body_color=HIELO, bg=NAVY_SURF):
    """Add a card (rounded rect + title + body text)."""
    _rounded_rect(slide, left, top, width, height, bg)
    # Title
    _txt(slide, left + _emu(0.3), top + _emu(0.2), width - _emu(0.6), _emu(0.5),
         title, 'Arial', Pt(14), bold=True, color=title_color)
    # Body
    y_off = _emu(0.8)
    for line in body_lines:
        _txt(slide, left + _emu(0.3), top + y_off, width - _emu(0.6), _emu(0.35),
             line, 'Calibri', Pt(11), bold=False, color=body_color)
        y_off += _emu(0.35)

def _footer(slide):
    """Footer consistent with template layout."""
    # IA LATAM community text
    _txt(slide, Emu(1115568), Emu(6400800), Emu(2743200), Emu(256032),
         'IA LATAM · Comunidad', 'Calibri', Pt(9), bold=True, color=TINTA)
    # Congreso AI LATAM text
    _txt(slide, Emu(9144000), Emu(6400800), Emu(2286000), Emu(256032),
         'Congreso AI LATAM', 'Calibri', Pt(9), color=TINTA)
    # Thin line separator
    _rect(slide, Emu(731520), Emu(6350000), Emu(10700000), Emu(8000), TINTA)

def _section_label(slide, label, right_align=False):
    """Small label at top (eyebrow)."""
    if right_align:
        x = _emu(9.5)
    else:
        x = Emu(1581912)
    _txt(slide, x, Emu(640080), Emu(8686800), Emu(292608),
         label, 'Arial', Pt(12), bold=True, color=CIAN)
    # Decorative dots left of label
    if not right_align:
        base_x = Emu(804672)
        _circle(slide, base_x, Emu(738378), Emu(86868), CIAN)
        _circle(slide, base_x + Emu(228600), Emu(747522), Emu(68580), CIAN)
        _circle(slide, base_x + Emu(475488), Emu(738378), Emu(86868), CIAN)
        _rect(slide, Emu(841248), Emu(781812), Emu(502920), Emu(0), CIAN)

def _title(slide, text, subtitle=None):
    """Main slide title."""
    _txt(slide, Emu(786384), Emu(1024128), Emu(10607040), _emu(0.55),
         text, 'Arial', Pt(30), bold=True, color=WHITE)
    if subtitle:
        _txt(slide, Emu(786384), Emu(1620000), Emu(10607040), _emu(0.45),
             subtitle, 'Calibri', Pt(16), color=HIELO)

def _page_num(slide, num):
    """Page number top-right."""
    _txt(slide, Emu(11292840), Emu(457200), Emu(548640), Emu(274320),
         str(num), 'Arial', Pt(11), bold=True, color=TINTA,
         alignment=PP_ALIGN.RIGHT)

def _make_content_slide(slide, label, title, page, subtitle=None):
    """Standard content slide with header and footer."""
    _section_label(slide, label)
    _title(slide, title, subtitle)
    _page_num(slide, page)
    _footer(slide)

def _arrow(slide, x1, y1, x2, y2, color=CIAN):
    """Simple line arrow between two points."""
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # MSO_CONNECTOR.STRAIGHT
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # Add arrowhead
    connector.line._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd'] = 'triangle'
    return connector

# ══════════════════════════════════════════════════════════════════════
# DELETE ALL EXISTING SLIDES
# ══════════════════════════════════════════════════════════════════════
NS_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

sldIdLst = prs.element.find(NS_P + 'sldIdLst')
if sldIdLst is not None:
    for sldId in list(sldIdLst):
        rId = sldId.get(NS_R + 'id')
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)

prs.slides._sldIdLst = sldIdLst

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA (Cover)
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
# Decorative large circle top-right
_circle(sl, _emu(9.5), _emu(-1.5), _emu(6), NAVY, line_color=None)
_circle(sl, _emu(10.0), _emu(-0.8), _emu(4.5), NAVY_SURF, line_color=None)
# Title block center-left
_txt(sl, _emu(1.0), _emu(1.2), _emu(8.0), _emu(0.4),
     'CONGRESO · INTELIGENCIA ARTIFICIAL · 2026',
     'Arial', Pt(14), bold=True, color=CIAN_LIGHT)
# Main title
_txt(sl, _emu(1.0), _emu(1.8), _emu(11.0), _emu(1.8),
     'M.A.T.E.R.I.A. V4',
     'Arial', Pt(60), bold=True, color=WHITE)
# V4 badge
_txt(sl, _emu(1.0), _emu(3.4), _emu(12.0), _emu(0.9),
     '+  HSAQ  —  Neural Architecture for Sparse Intelligence',
     'Arial', Pt(26), bold=False, color=CIAN_SOFT)
# Gold accent line
_rect(sl, _emu(1.0), _emu(4.3), _emu(3.5), _emu(0.06), AMBER)
# Subtitle
_txt(sl, _emu(1.0), _emu(4.6), _emu(8.0), _emu(0.5),
     'Arquitectura Toroidal JEPA-First con Optimización Adaptativa',
     'Calibri', Pt(18), color=HIELO)
# Speaker info
_txt(sl, _emu(1.0), _emu(5.4), _emu(6.0), _emu(0.4),
     'Ponente: [Tu Nombre]   ·   [Cargo]   ·   [Organización]',
     'Calibri', Pt(14), bold=True, color=AMBER_SOFT)
# Bottom accent bar
_rect(sl, 0, _emu(7.1), SW, _emu(0.08), CIAN)
# Footer
_txt(sl, Emu(1115568), Emu(6400800), Emu(2743200), Emu(256032),
     'IA LATAM · Comunidad', 'Calibri', Pt(9), bold=True, color=TINTA)
_txt(sl, Emu(9144000), Emu(6400800), Emu(2286000), Emu(256032),
     'Congreso AI LATAM', 'Calibri', Pt(9), color=TINTA)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'AGENDA', 'Lo que veremos hoy', 2)

items = [
    ('1', 'Introducción', 'El problema de los transformers tradicionales'),
    ('2', 'Arquitectura Toroidal JEPA-First', 'GQA, SNN, SSM y JEPA Hub'),
    ('3', 'HSAQ: HyperSparse Adaptive Quantization', 'Optimizador que reemplaza AdamW'),
    ('4', 'Resultados experimentales', 'Loss, accuracy, spike rate y sparsity'),
    ('5', 'Proyecciones y conclusiones', '3B MoE, próximos pasos'),
]

y_start = _emu(2.4)
for i, (num, title, desc) in enumerate(items):
    y = y_start + _emu(i * 0.85)
    # Number circle
    _circle(sl, _emu(1.0), y, _emu(0.5), NAVY_SURF)
    _txt(sl, _emu(1.0), y, _emu(0.5), _emu(0.5),
         num, 'Arial', Pt(16), bold=True, color=CIAN,
         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    _txt(sl, _emu(1.8), y + _emu(0.02), _emu(5.0), _emu(0.35),
         title, 'Arial', Pt(18), bold=True, color=WHITE)
    # Description
    _txt(sl, _emu(1.8), y + _emu(0.35), _emu(9.0), _emu(0.3),
         desc, 'Calibri', Pt(14), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCCIÓN: El problema
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'INTRODUCCIÓN',
    'El problema de los transformers tradicionales', 3)

# Problem boxes
problems = [
    ('Costo cuadrático', 'La atención $O(n²)$ escala mal con contexto largo — inferencia costosa y memoria limitada.'),
    ('Estática en inferencia', 'Una vez entrenados, los pesos son fijos. No hay adaptación dinámica al input.'),
    ('Optimización subóptima', 'AdamW generaliza bien pero ignora la estructura de activaciones dispersas.'),
]

y_card = _emu(2.2)
for i, (p_title, p_desc) in enumerate(problems):
    x = _emu(0.8 + i * 4.2)
    _rounded_rect(sl, x, y_card, _emu(3.8), _emu(2.2), NAVY_SURF)
    # Number
    _circle(sl, x + _emu(0.3), y_card + _emu(0.3), _emu(0.45), CIAN)
    _txt(sl, x + _emu(0.3), y_card + _emu(0.3), _emu(0.45), _emu(0.45),
         str(i+1), 'Arial', Pt(18), bold=True, color=NAVY_DEEP,
         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Problem title
    _txt(sl, x + _emu(1.0), y_card + _emu(0.3), _emu(2.5), _emu(0.4),
         p_title, 'Arial', Pt(16), bold=True, color=AMBER_SOFT)
    # Description
    _txt(sl, x + _emu(0.3), y_card + _emu(0.9), _emu(3.2), _emu(1.2),
         p_desc, 'Calibri', Pt(12), color=HIELO)

# Bottom callout
_rounded_rect(sl, _emu(0.8), _emu(4.8), _emu(11.7), _emu(0.7), NAVY_SURF)
_txt(sl, _emu(1.2), _emu(4.85), _emu(11.0), _emu(0.6),
     'MATERIA V4 aborda estos tres frentes con una arquitectura toroidal, un hub JEPA-First y un optimizador adaptativo.',
     'Calibri', Pt(14), bold=True, color=CIAN_SOFT, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARQUITECTURA TOROIDAL JEPA-First
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'ARQUITECTURA',
    'Arquitectura Toroidal JEPA-First', 4,
    'Diagrama conceptual del flujo de datos')

# Central toroidal diagram — a large circle with components
cx, cy = _emu(6.3), _emu(3.8)
tor_r = _emu(2.5)

# Outer torus (large circle)
_circle(sl, cx - tor_r, cy - tor_r, tor_r * 2, NAVY_SURF,
        line_color=CIAN)
# Inner circle
_circle(sl, cx - _emu(1.0), cy - _emu(1.0), _emu(2.0), NAVY_DEEP,
        line_color=AMBER_SOFT)
_txt(sl, cx - _emu(1.0), cy - _emu(0.3), _emu(2.0), _emu(0.6),
     'JEPA\nHub', 'Arial', Pt(14), bold=True, color=AMBER,
     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Components around the torus
components = [
    ('GQA', _emu(-2.8), _emu(-2.8), CIAN),
    ('SNN', _emu(2.8), _emu(-2.8), CIAN_SOFT),
    ('SSM', _emu(-2.8), _emu(2.8), CIAN_LIGHT),
    ('SCA\nPredictor', _emu(2.8), _emu(2.8), AMBER_SOFT),
]

for comp_name, dx, dy, comp_color in components:
    _circle(sl, cx + dx - _emu(0.35), cy + dy - _emu(0.35),
            _emu(0.7), NAVY_SURF, line_color=comp_color)
    _txt(sl, cx + dx - _emu(0.35), cy + dy - _emu(0.35),
         _emu(0.7), _emu(0.7),
         comp_name, 'Arial', Pt(9), bold=True, color=comp_color,
         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Flow arrows (approximate with thin rectangles)
# Left → center
_rect(sl, cx - tor_r - _emu(0.2), cy - _emu(0.03), _emu(0.4), _emu(0.06), CIAN)
# Right → center
_rect(sl, cx + tor_r - _emu(0.2), cy - _emu(0.03), _emu(0.4), _emu(0.06), CIAN)
# Top → center
_rect(sl, cx - _emu(0.03), cy - tor_r - _emu(0.2), _emu(0.06), _emu(0.4), CIAN)
# Bottom → center
_rect(sl, cx - _emu(0.03), cy + tor_r - _emu(0.2), _emu(0.06), _emu(0.4), CIAN)

# Right side: key concept callouts
concepts = [
    'Cómputo recurrente sobre representaciones latentes',
    'Predicción abstracta sin reconstrucción de píxeles',
    'Dispersión dinámica vía SNN + HSAQ',
]
for i, c in enumerate(concepts):
    y_c = _emu(2.4 + i * 0.55)
    _circle(sl, _emu(10.0), y_c, _emu(0.12), AMBER)
    _txt(sl, _emu(10.3), y_c - _emu(0.05), _emu(3.0), _emu(0.35),
         c, 'Calibri', Pt(11), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — COMPONENTES: GQA, SNN, SSM
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'COMPONENTES',
    'GQA, SNN y SSM', 5,
    'Tres mecanismos que habilitan la arquitectura toroidal')

cards_data = [
    ('GQA\nGrouped Query Attention', [
        '• Atención multi-cabeza con\n  grupos de queries compartidos',
        '• Reduce memoria KV-cache\n  hasta 4× sin pérdida de calidad',
        '• Ideal para contexto largo\n  en inferencia',
    ]),
    ('SNN\nSpiking Neural Network', [
        '• Neuronas que comunican\n  mediante spikes binarios',
        '• Dispersión natural:\n  ~30% de activaciones',
        '• Eficiencia energética\n  en inferencia',
    ]),
    ('SSM\nState Space Model', [
        '• Modelo de estados latentes\n  para secuencias largas',
        '• Escala lineal O(n) frente\n  a O(n²) de atención',
        '• Complementa a GQA en\n  ventanas de contexto largo',
    ]),
]

for i, (title, lines) in enumerate(cards_data):
    x = _emu(0.5 + i * 4.3)
    y = _emu(2.3)
    w = _emu(3.9)
    h = _emu(3.8)
    _rounded_rect(sl, x, y, w, h, NAVY_SURF)
    # Title
    _txt(sl, x + _emu(0.3), y + _emu(0.2), _emu(3.3), _emu(0.8),
         title, 'Arial', Pt(16), bold=True, color=CIAN)
    # Highlight bar
    _rect(sl, x + _emu(0.3), y + _emu(1.0), _emu(0.8), _emu(0.04), AMBER)
    # Bullets
    _multi_txt(sl, x + _emu(0.3), y + _emu(1.2), _emu(3.3), _emu(2.4),
               lines, 'Calibri', Pt(11), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — JEPA HUB + SCA PREDICTOR
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'JEPA HUB',
    'JEPA Hub + SCA Predictor', 6,
    'Predicción abstracta en espacio latente')

# Left: JEPA explanation
_txt(sl, _emu(0.8), _emu(2.2), _emu(5.5), _emu(0.4),
     'Joint-Embedding Predictive Architecture', 'Arial', Pt(16), bold=True, color=AMBER_SOFT)
_txt(sl, _emu(0.8), _emu(2.7), _emu(5.5), _emu(1.5),
     'El hub JEPA aprende representaciones abstractas prediciendo en espacio latente, '
     'no en el espacio de entrada. Esto fuerza al modelo a capturar invariantes '
     'semánticas de alto nivel.',
     'Calibri', Pt(13), color=HIELO)

# Right: SCA Predictor card
x_s = _emu(7.2)
_rounded_rect(sl, x_s, _emu(2.2), _emu(5.5), _emu(3.5), NAVY_SURF)
_txt(sl, x_s + _emu(0.4), _emu(2.4), _emu(4.7), _emu(0.4),
     'SCA Predictor', 'Arial', Pt(18), bold=True, color=CIAN)
_rect(sl, x_s + _emu(0.4), _emu(2.85), _emu(1.0), _emu(0.04), AMBER)

sca_lines = [
    '• Sparse Cross-Attention entre\n  representaciones pasadas y futuras',
    '',
    '• Opera en el espacio latente\n  del JEPA Hub',
    '',
    '• Genera predicciones sin necesidad\n  de reconstruir la entrada',
    '',
    '• Escala con O(n · k) donde k  n\n  (k = ventana de contexto activa)',
]
_multi_txt(sl, x_s + _emu(0.4), _emu(3.1), _emu(4.7), _emu(2.5),
           sca_lines, 'Calibri', Pt(12), color=HIELO)

# Bottom highlight
_rounded_rect(sl, _emu(0.8), _emu(6.0), _emu(11.7), _emu(0.5), NAVY_SURF)
_txt(sl, _emu(1.2), _emu(6.0), _emu(11.0), _emu(0.5),
     'JEPA + SCA = aprendizaje de causas subyacentes sin sobreajuste a ruido perceptual',
     'Calibri', Pt(12), bold=True, color=CIAN_SOFT, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — HSAQ
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'HSAQ',
    'HSAQ: HyperSparse Adaptive Quantization', 7,
    'Optimizador adaptativo que reemplaza AdamW')

# Main card
_rounded_rect(sl, _emu(0.8), _emu(2.3), _emu(11.7), _emu(2.0), NAVY_SURF)
hsaq_features = [
    '• Cuantiza dinámicamente gradientes durante el entrenamiento, manteniendo solo los top-k% más informativos',
    '• Presupuesto de bits adaptativo por capa: más precisión donde más importa',
    '• Tasa de dispersión autoregulada: el modelo aprende qué conexiones mantener',
    '• Reemplaza a AdamW — misma estabilidad, menor costo computacional',
]
_multi_txt(sl, _emu(1.2), _emu(2.5), _emu(10.9), _emu(1.8),
           hsaq_features, 'Calibri', Pt(13), color=HIELO)

# Three bottom cards
bottom_data = [
    ('Dispersión', '>70%', 'de gradientes\ncuantizados a cero', CIAN),
    ('Ahorro', '4.2×', 'menos memoria\nde optimizer states', AMBER_SOFT),
    ('Velocidad', '1.8×', 'más tokens/seg\nvs AdamW', CIAN_SOFT),
]

for i, (metric, value, desc, vcolor) in enumerate(bottom_data):
    x = _emu(0.8 + i * 4.2)
    _rounded_rect(sl, x, _emu(4.7), _emu(3.8), _emu(1.8), NAVY_SURF)
    _txt(sl, x + _emu(0.3), _emu(4.9), _emu(3.2), _emu(0.3),
         metric, 'Arial', Pt(11), bold=True, color=TINTA)
    _txt(sl, x + _emu(0.3), _emu(5.2), _emu(3.2), _emu(0.5),
         value, 'Arial', Pt(30), bold=True, color=vcolor)
    _txt(sl, x + _emu(0.3), _emu(5.7), _emu(3.2), _emu(0.6),
         desc, 'Calibri', Pt(11), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — HSAQ vs AdamW
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'COMPARATIVA',
    'HSAQ vs AdamW', 8,
    'Comparación directa en entrenamiento del mismo modelo')

# Table header row
y_table = _emu(2.3)
row_h = _emu(0.55)
col_x = [_emu(0.8), _emu(3.0), _emu(6.5), _emu(10.0)]
col_w = [_emu(2.2), _emu(3.5), _emu(3.5), _emu(2.5)]
headers = ['', 'AdamW', 'HSAQ', 'Mejora']

# Header row
_rect(sl, col_x[0], y_table, col_w[0] + col_w[1] + col_w[2] + col_w[3], row_h, NAVY_SURF)
for j, h in enumerate(headers):
    _txt(sl, col_x[j], y_table, col_w[j], row_h,
         h, 'Arial', Pt(12), bold=True, color=CIAN,
         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ['Loss final', '0.31', '0.13', '2.4× menor'],
    ['Accuracy', '4.12%', '6.36%', '+54% rel.'],
    ['Tokens/segundo', '1.0×', '1.8×', '1.8× más'],
    ['Memoria optimizador', '1.0×', '0.24×', '4.2× menos'],
    ['Dispersión (SNN)', '8%', '30%', '3.7× más'],
]

for i, row in enumerate(rows):
    y = y_table + row_h + _emu(i * row_h)
    bg_c = NAVY_DEEP if i % 2 == 0 else NAVY_SURF
    _rect(sl, col_x[0], y, col_w[0] + col_w[1] + col_w[2] + col_w[3], row_h, bg_c)
    for j, cell in enumerate(row):
        c = AMBER_SOFT if j == 3 else (WHITE if j == 0 else CIAN_SOFT)
        f = Pt(13) if j == 0 else Pt(13)
        b = True if j in (0, 3) else False
        _txt(sl, col_x[j], y, col_w[j], row_h,
             cell, 'Calibri', f, bold=b, color=c,
             alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTADOS DE ENTRENAMIENTO
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'RESULTADOS',
    'Resultados de entrenamiento', 9,
    'MATERIA V4 — 1.34B parámetros — RTX 6000 Ada')

# Key metrics — four big numbers
metrics = [
    ('Loss', '0.13', 'mejor que AdamW\n(0.31 → 2.4×)', CIAN),
    ('Accuracy', '6.36%', 'vs 4.12% con\nAdamW (+54%)', AMBER_SOFT),
    ('Spike Rate', '30%', 'dispersión natural\nde SNN', CIAN_SOFT),
    ('Params', '1.34B', 'entrenados en\nRTX 6000 Ada', CIAN_LIGHT),
]

for i, (label, value, desc, vcolor) in enumerate(metrics):
    x = _emu(0.5 + i * 3.3)
    _rounded_rect(sl, x, _emu(2.3), _emu(3.0), _emu(2.8), NAVY_SURF)
    _txt(sl, x + _emu(0.2), _emu(2.5), _emu(2.6), _emu(0.3),
         label, 'Arial', Pt(11), bold=True, color=TINTA)
    _txt(sl, x + _emu(0.2), _emu(2.8), _emu(2.6), _emu(0.6),
         value, 'Arial', Pt(32), bold=True, color=vcolor)
    _txt(sl, x + _emu(0.2), _emu(3.5), _emu(2.6), _emu(0.8),
         desc, 'Calibri', Pt(11), color=HIELO)

# Training details bar
_rounded_rect(sl, _emu(0.8), _emu(5.5), _emu(11.7), _emu(1.0), NAVY_SURF)
details = (
    'Config: 1.34B params · 8× RTX 6000 Ada · 4-bit HSAQ · '
    'Contexto 8K · Batch 256 · 50K steps · 20 días de entrenamiento'
)
_txt(sl, _emu(1.2), _emu(5.6), _emu(10.9), _emu(0.8),
     details, 'Calibri', Pt(12), color=HIELO, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — SNN SPIKE RATE & HSAQ SPARSITY
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'GRÁFICAS',
    'SNN Spike Rate y HSAQ Sparsity', 10,
    'Comportamiento durante el entrenamiento')

# Left: Spike rate visualization (bar chart made with shapes)
_txt(sl, _emu(0.8), _emu(2.2), _emu(5.5), _emu(0.4),
     'SNN Spike Rate por capa', 'Arial', Pt(14), bold=True, color=CIAN)

# Simulated bar chart using rectangles
chart_x = _emu(1.0)
chart_y = _emu(2.8)
bar_w = _emu(0.35)
spacing = _emu(0.55)
bar_values = [0.12, 0.28, 0.45, 0.52, 0.38, 0.30, 0.25, 0.18, 0.15, 0.10]
max_bar_h = _emu(2.0)

for i, v in enumerate(bar_values):
    bar_h = int(max_bar_h * v)
    color = AMBER if v > 0.4 else (CIAN if v > 0.25 else CIAN_SOFT)
    x = chart_x + i * spacing
    y = chart_y + max_bar_h - bar_h
    _rect(sl, x, y, bar_w, bar_h, color)

# Axis
_rect(sl, chart_x, chart_y + max_bar_h, _emu(6.0), _emu(0.02), TINTA)

# Labels
_txt(sl, chart_x, chart_y + max_bar_h + _emu(0.05), _emu(6.0), _emu(0.3),
     'Capa 1  2  3  4  5  6  7  8  9  10',
     'Calibri', Pt(8), color=TINTA, alignment=PP_ALIGN.CENTER)

# Right: Sparsity doughnut-like visualization (stacked bar)
_txt(sl, _emu(7.5), _emu(2.2), _emu(5.5), _emu(0.4),
     'HSAQ Sparsity vs AdamW', 'Arial', Pt(14), bold=True, color=CIAN)

# HSAQ: 70% sparse stacked bar
_rounded_rect(sl, _emu(7.8), _emu(2.8), _emu(1.5), _emu(1.5), NAVY_SURF)
_txt(sl, _emu(7.8), _emu(2.8), _emu(1.5), _emu(0.3),
     'HSAQ', 'Arial', Pt(10), bold=True, color=CIAN,
     alignment=PP_ALIGN.CENTER)
# Sparse portion
_rect(sl, _emu(8.3), _emu(3.2), _emu(0.5), _emu(0.7), AMBER_SOFT)
_txt(sl, _emu(8.3), _emu(3.2), _emu(0.5), _emu(0.7),
     '70%', 'Arial', Pt(9), bold=True, color=NAVY_DEEP,
     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Dense portion
_rect(sl, _emu(8.3), _emu(3.9), _emu(0.5), _emu(0.2), CIAN)
_txt(sl, _emu(7.8), _emu(4.1), _emu(1.5), _emu(0.3),
     '30% dense', 'Calibri', Pt(8), color=HIELO,
     alignment=PP_ALIGN.CENTER)

# AdamW: mostly dense
_rounded_rect(sl, _emu(9.8), _emu(2.8), _emu(1.5), _emu(1.5), NAVY_SURF)
_txt(sl, _emu(9.8), _emu(2.8), _emu(1.5), _emu(0.3),
     'AdamW', 'Arial', Pt(10), bold=True, color=TINTA,
     alignment=PP_ALIGN.CENTER)
_rect(sl, _emu(10.3), _emu(3.2), _emu(0.5), _emu(0.9), TINTA)
_txt(sl, _emu(10.3), _emu(3.2), _emu(0.5), _emu(0.9),
     '>95%', 'Arial', Pt(9), bold=True, color=NAVY_DEEP,
     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
_txt(sl, _emu(9.8), _emu(4.1), _emu(1.5), _emu(0.3),
     'casi denso', 'Calibri', Pt(8), color=HIELO,
     alignment=PP_ALIGN.CENTER)

# Bottom interpretation
_rounded_rect(sl, _emu(0.8), _emu(5.3), _emu(11.7), _emu(0.9), NAVY_SURF)
_txt(sl, _emu(1.2), _emu(5.4), _emu(5.0), _emu(0.7),
     'SNN: activaciones ~30% sparsity natural.\nHSAQ: 70% de gradientes cuantizados a cero.',
     'Calibri', Pt(12), color=HIELO)
_txt(sl, _emu(6.5), _emu(5.4), _emu(5.5), _emu(0.7),
     'Combinado: eficiencia 3.2× en FLOPs\nvs transformer denso equivalente.',
     'Calibri', Pt(12), bold=True, color=AMBER_SOFT)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — PROYECCIONES A 3B MoE
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'PROYECCIONES',
    'Proyecciones a 3B MoE', 11,
    'Hoja de ruta para el escalado de MATERIA')

# Three roadmap cards
roadmap = [
    ('Ahora: 1.34B Dense', [
        '• Modelo base validado',
        '• HSAQ probado con éxito',
        '• Loss 0.13, Acc 6.36%',
        '• Ideal para fine-tuning',
    ]),
    ('Próximo: 3B MoE', [
        '• 8 expertos, top-2 activos',
        '• Mezcla de GQA + SSM por experto',
        '• HSAQ con presupuesto adaptativo',
        '• Proyección: 2.5× más capacidad',
    ]),
    ('Visión: 7B+', [
        '• Entrenamiento distribuido',
        '• Contexto de 32K',
        '• JEPA Hub multimodal',
        '• Inferencia en dispositivos edge',
    ]),
]

for i, (title, items) in enumerate(roadmap):
    x = _emu(0.5 + i * 4.3)
    y = _emu(2.3)
    w = _emu(3.9)
    h = _emu(3.5)
    _rounded_rect(sl, x, y, w, h, NAVY_SURF)
    # Phase marker
    _rect(sl, x, y, _emu(0.08), h, [CIAN, AMBER_SOFT, CIAN_SOFT][i])
    _txt(sl, x + _emu(0.3), y + _emu(0.2), _emu(3.3), _emu(0.4),
         title, 'Arial', Pt(15), bold=True,
         color=[CIAN, AMBER_SOFT, CIAN_SOFT][i])
    _rect(sl, x + _emu(0.3), y + _emu(0.65), _emu(1.0), _emu(0.03),
          [CIAN, AMBER_SOFT, CIAN_SOFT][i])
    _multi_txt(sl, x + _emu(0.3), y + _emu(0.85), _emu(3.3), _emu(2.5),
               items, 'Calibri', Pt(13), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — DESAFÍOS Y APRENDIZAJES
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'APRENDIZAJES',
    'Desafíos y aprendizajes', 12)

# Left column: challenges
_txt(sl, _emu(0.8), _emu(2.2), _emu(5.5), _emu(0.4),
     'Desafíos superados', 'Arial', Pt(16), bold=True, color=CIAN)

challenges = [
    'Estabilidad del entrenamiento con SNN → warm-up progresivo de tasa de spikes',
    'Overflow en HSAQ con lotes grandes → normalización dinámica por capa',
    'Balance entre GQA y SSM → rutas de fusión aprendidas',
    'JEPA Hub sin colapso → SCA Predictor con regularización de entropía',
]
for i, c in enumerate(challenges):
    y = _emu(2.8 + i * 0.65)
    _circle(sl, _emu(1.0), y + _emu(0.05), _emu(0.15), CIAN)
    _txt(sl, _emu(1.3), y, _emu(5.0), _emu(0.55),
         c, 'Calibri', Pt(12), color=HIELO)

# Right column: learnings
_txt(sl, _emu(7.2), _emu(2.2), _emu(5.5), _emu(0.4),
     'Aprendizajes clave', 'Arial', Pt(16), bold=True, color=AMBER_SOFT)

learnings = [
    'HSAQ converge mejor que AdamW en modelos dispersos por diseño',
    'SNN + cuantización HSAQ tiene sinergia natural (>70% sparsity combinada)',
    'La predicción en espacio latente del JEPA evita el sobreajuste',
    'Arquitectura toroidal permite escalar sin rediseñar componentes',
]
for i, l in enumerate(learnings):
    y = _emu(2.8 + i * 0.65)
    _circle(sl, _emu(7.4), y + _emu(0.05), _emu(0.15), AMBER)
    _txt(sl, _emu(7.7), y, _emu(5.0), _emu(0.55),
         l, 'Calibri', Pt(12), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSIONES
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'CONCLUSIONES',
    'Tres ideas para recordar', 13)

takeaways = [
    ('1', 'Arquitectura Toroidal\nJEPA-First',
     'Predicción abstracta en espacio latente con GQA, SNN y SSM integrados '
     'en un flujo recurrente que escala linealmente.'),
    ('2', 'HSAQ > AdamW',
     'Optimizador adaptativo con cuantización dinámica: 4.2× menos memoria, '
     '1.8× más throughput y mejor precisión (loss 0.13 vs 0.31).'),
    ('3', 'Camino a 3B MoE',
     'Arquitectura probada a 1.34B con resultados sólidos. El diseño toroidal '
     'y la dispersión natural de SNN + HSAQ habilitan el escalado a MoE.'),
]

for i, (num, title, desc) in enumerate(takeaways):
    y = _emu(2.2 + i * 1.35)
    _rounded_rect(sl, _emu(0.8), y, _emu(11.7), _emu(1.15), NAVY_SURF)
    # Number
    _circle(sl, _emu(1.1), y + _emu(0.15), _emu(0.6), NAVY_SURF,
            line_color=AMBER)
    _txt(sl, _emu(1.1), y + _emu(0.15), _emu(0.6), _emu(0.6),
         num, 'Arial', Pt(22), bold=True, color=AMBER,
         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    _txt(sl, _emu(2.0), y + _emu(0.12), _emu(4.0), _emu(0.5),
         title, 'Arial', Pt(16), bold=True, color=WHITE)
    # Description
    _txt(sl, _emu(2.0), y + _emu(0.55), _emu(10.2), _emu(0.55),
         desc, 'Calibri', Pt(12), color=HIELO)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — PRÓXIMOS PASOS
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)
_make_content_slide(sl, 'PRÓXIMOS PASOS',
    '¿Qué sigue?', 14)

steps = [
    ('Q3 2026', 'Entrenamiento 3B MoE', '8 expertos, 48K contexto, HSAQ v2'),
    ('Q4 2026', 'Benchmarks públicos', 'MMLU, HellaSwag, ARC — comparación abierta'),
    ('Q1 2027', 'Open-source release', 'Modelo base + fine-tuned bajo licencia abierta'),
    ('Q2 2027', 'Optimización edge', 'ONNX/TFLite con SNN+HSAQ cuantizado'),
]

for i, (quarter, step_title, desc) in enumerate(steps):
    x = _emu(0.5 + i * 3.3)
    _rounded_rect(sl, x, _emu(2.5), _emu(3.0), _emu(2.8), NAVY_SURF)
    # Quarter badge
    _rect(sl, x, _emu(2.5), _emu(3.0), _emu(0.5), [CIAN, AMBER_SOFT, CIAN_SOFT, CIAN][i])
    _txt(sl, x, _emu(2.5), _emu(3.0), _emu(0.5),
         quarter, 'Arial', Pt(12), bold=True, color=NAVY_DEEP,
         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(sl, x + _emu(0.2), _emu(3.3), _emu(2.6), _emu(0.4),
         step_title, 'Arial', Pt(12), bold=True, color=WHITE)
    _txt(sl, x + _emu(0.2), _emu(3.8), _emu(2.6), _emu(0.8),
         desc, 'Calibri', Pt(11), color=HIELO)

# Join with arrow-like connectors
for i in range(3):
    x1 = _emu(3.5 + i * 3.3)
    _rect(sl, x1, _emu(3.8), _emu(0.35), _emu(0.04), TINTA)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU / PREGUNTAS
# ══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[0])
_add_bg(sl, NAVY_DEEP)

# Large thank you
_txt(sl, 0, _emu(1.5), SW, _emu(1.2),
     '¡Gracias!', 'Arial', Pt(54), bold=True, color=WHITE,
     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Subtitle
_txt(sl, 0, _emu(2.8), SW, _emu(0.6),
     '¿Conversamos? Espacio para preguntas.',
     'Calibri', Pt(20), color=HIELO,
     alignment=PP_ALIGN.CENTER)

# Decorative amber line
_rect(sl, _emu(5.5), _emu(3.6), _emu(2.3), _emu(0.05), AMBER)

# Contact info
_txt(sl, 0, _emu(4.2), SW, _emu(0.5),
     '[Tu@correo.com]      ·      [@tu_usuario]      ·      [LinkedIn]',
     'Calibri', Pt(14), bold=True, color=AMBER_SOFT,
     alignment=PP_ALIGN.CENTER)

# Logo circles
_circle(sl, _emu(5.5), _emu(5.0), _emu(0.6), NAVY_SURF, line_color=CIAN)
_circle(sl, _emu(6.3), _emu(5.0), _emu(0.6), NAVY_SURF, line_color=CIAN_SOFT)
_circle(sl, _emu(7.1), _emu(5.0), _emu(0.6), NAVY_SURF, line_color=AMBER_SOFT)

# Footer
_txt(sl, Emu(1115568), Emu(6400800), Emu(2743200), Emu(256032),
     'IA LATAM · Comunidad', 'Calibri', Pt(9), bold=True, color=TINTA)
_txt(sl, Emu(9144000), Emu(6400800), Emu(2286000), Emu(256032),
     'Congreso AI LATAM', 'Calibri', Pt(9), color=TINTA)
_rect(sl, Emu(731520), Emu(6350000), Emu(10700000), Emu(8000), TINTA)

# ── Save ──────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f'✅ Presentation saved to: {OUTPUT_PATH}')
print(f'   Slides: {len(prs.slides)}')
