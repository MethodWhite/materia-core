"""
Genera PPTX presentación MATERIA V4 + HSAQ para Congreso AI LATAM
Basado en la plantilla Plantilla_Congreso_AI_LATAM.pptx
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

TEMPLATE = "/home/methodwhite/Documentos/IA LATAM/Plantilla/Plantilla_Congreso_AI_LATAM.pptx"
OUTPUT   = "/home/methodwhite/Documentos/IA LATAM/M.A.T.E.R.I.A./MATERIA_HSAQ_Conferencia_AI_LATAM.pptx"

# ── Content mapping: (slide_idx, text_replacements) ──
# Each tuple: (partial match, new text)
# We'll replace text in text frames that contain the match string.

CONTENT = {
    # Slide 1 - Portada
    "CONGRESO · INTELIGENCIA ARTIFICIAL · 2026": "M.A.T.E.R.I.A. V4 · ARQUITECTURA HSAQ · 2026",
    "Congreso\nAI LATAM": "M.A.T.E.R.I.A. V4\n+ HSAQ",
    "El encuentro de la comunidad de Inteligencia Artificial de Latinoamérica.": "Arquitectura de Cuantización Adaptativa Hiperdispersa para modelos de próxima generación.",
    "Fecha del evento": "30 de Julio, 2026",
    "Ciudad · Sede": "Online · LATAM",
    "ialatam.com": "github.com/MethodWhite/MATERIA",

    # Slide 2 - Charla magistral
    "CHARLA MAGISTRAL · TRACK DE IA": "M.A.T.E.R.I.A. V4 · ARQUITECTURA HSAQ",
    "Título de tu charla en\nuna o dos líneas": "M.A.T.E.R.I.A. V4\nArquitectura HSAQ",
    "Tu Nombre Apellido   ·   Cargo · Organización": "MethodWhite   ·   Investigador · M.A.T.E.R.I.A. Research",

    # Slide 4 - Agenda
    "Lo que veremos hoy": "Lo que veremos hoy",
    # We'll modify the bullet list items

    # Slide 5 - Sobre el ponente
    "SOBRE EL PONENTE": "M.A.T.E.R.I.A.",
    "Quién te acompaña hoy": "¿Qué es M.A.T.E.R.I.A.?",

    # Slide 6 - Contenido (genérico)
    "Título de la lámina de contenido": "¿Qué es HSAQ?",

    # Slide 7 - Comparativa
    "Dos columnas de contenido": "HSAQ vs TurboQuant",

    # Slide 10 - Conceptos
    "Tres ideas en tarjetas": "Tres pilares de HSAQ",

    # Slide 11 - Pilares
    "Cuatro bloques en cuadrícula": "Aplicación por capas",

    # Slide 12 - Proceso
    "Un proceso en cuatro pasos": "Pipeline de entrenamiento",

    # Slide 13 - Comparación
    "Enfoque tradicional vs. con IA": "Sin HSAQ vs Con HSAQ",

    # Slide 14 - Cifras
    "Los datos que respaldan tu mensaje": "Métricas y resultados",

    # Slide 15 - Cita
    "Autor de la cita  ·  Cargo / referencia": "MethodWhite  ·  M.A.T.E.R.I.A. Research 2026",

    # Slide 18 - Arquitectura
    "Diagrama de flujo o arquitectura": "Arquitectura M.A.T.E.R.I.A. V4",

    # Slide 20 - Trayectoria
    "Línea de tiempo del proyecto": "Evolución de HSAQ",

    # Slide 22 - Para llevar
    "Tres ideas para recordar": "Conclusiones",

    # Slide 23 - Gracias
    "¿Conversamos? Espacio para preguntas.": "¿Conversamos? methodwhite@github.com",
    "tu@correo.com      ·      @tu_usuario": "methodwhite@github.com  ·  @MethodWhite",
}

CONTENT_DETAILS = {
    4: {  # Agenda
        "items": [
            "¿Qué es M.A.T.E.R.I.A. V4?",
            "Problema: Cuantización tradicional",
            "HSAQ: Cuantización adaptativa",
            "Arquitectura por capas",
            "Resultados y métricas",
            "Demostración en vivo",
        ]
    },
    5: {  # Sobre el ponente / Qué es MATERIA
        "text_blocks": [
            "M.A.T.E.R.I.A. es un modelo híbrido que integra:\n\n"
            "• Transformers con Flash Attention 2\n"
            "• LIF-SNN (neuronas pulsantes)\n"
            "• State Space Models (SSM)\n"
            "• JEPA (Joint Embedding Predictive Architecture)\n\n"
            "Todo entrenado con HSAQ: nuestro sistema de cuantización adaptativa que reemplaza a AdamW como optimizador."
        ]
    },
    6: {  # Qué es HSAQ
        "text_blocks": [
            "HyperSparse Adaptive Quantization\n\n"
            "HSAQ no es INT8 ni INT4. Es cuantización de activaciones mediante sparsity adaptativa dinámica:\n\n"
            "• kthvalue → umbral dinámico por batch\n"
            "• Máscara binaria {0, valor_original}\n"
            "• Sin calibración, sin datasets externos\n"
            "• Hardware-agnostic (CPU/GPU/TPU)"
        ]
    },
    7: {  # HSAQ vs TurboQuant
        "text_blocks": [
            "TurboQuant (Google):\n"
            "• Cuantización fija INT8\n"
            "• Requiere calibración post-entrenamiento\n"
            "• Solo GPU con soporte INT8\n"
            "• No adaptativo\n\n"
            "HSAQ (MATERIA):\n"
            "• Sparsity adaptativa por batch\n"
            "• Sin calibración\n"
            "• CPU/GPU/TPU\n"
            "• Umbral dinámico por capa"
        ]
    },
    10: {  # Tres pilares
        "text_blocks": [
            "ADAPTATIVO\nUmbral kthvalue\nse recalcula\nen cada batch",
            "POR CAPAS\nCada componente\ntiene su propia\nsparsity",
            "SIN ESTADO\nNo hay buffers\npersistentes\nentre batches"
        ]
    },
    11: {  # Aplicación por capas
        "text_blocks": [
            "Embedding: 5%",
            "Transformer (c/3): 8-15%",
            "SNN: 10%",
            "SSM: 5%",
            "JEPA: 5%"
        ]
    },
    12: {  # Pipeline
        "text_blocks": [
            "1. Forward pass\nEmbedding → HSAQ → Transformer → HSAQ → SNN → HSAQ → SSM → HSAQ → JEPA → HSAQ",
            "2. Sparsity escalonada\n0.05 → 0.15 progresivo\n(7 puntos estratégicos)",
            "3. Backward pass\nGradiente fluye solo por\nneuronas activas (STE)",
            "4. Weight update\nSGD Nesterov\n(momentum=0.9)"
        ]
    },
    13: {  # Sin HSAQ vs Con HSAQ
        "text_blocks": [
            "SIN HSAQ:\n"
            "• 100% de neuronas activas\n"
            "• Mayor riesgo de overfitting\n"
            "• AdamW: 8 bytes/param\n"
            "• SNN desactivado (spike=0)\n"
            "• 16.5M params máx",
            "CON HSAQ:\n"
            "• 70% de neuronas activas (regulariza)\n"
            "• Menos overfitting\n"
            "• SGD: 4 bytes/param (-50% VRAM)\n"
            "• SNN activo (spike=0.04)\n"
            "• 190M params (11.5× más)"
        ]
    },
    14: {  # Métricas
        "text_blocks": [
            "190M parámetros",
            "11.5× más params",
            "28.9% accuracy",
            "+48% mejora vs uniforme",
            "52% info preservada",
            "vs 0.7% sin escalonado"
        ]
    },
    15: {  # Cita
        "quote": "La Inteligencia Artificial no reemplaza el talento de la región, lo potencia. HSAQ democratiza el acceso a modelos grandes al reducir los requisitos de hardware.",
    },
    18: {  # Arquitectura
        "text_blocks": [
            "Input → Embedding → HSAQ(5%) → Block[0..9] → HSAQ(8-15%) → SNN → HSAQ(10%) → SSM → HSAQ(5%) → JEPA → HSAQ(5%) → Head"
        ]
    },
    20: {  # Evolución HSAQ
        "text_blocks": [
            "V1: Sparsity uniforme 30%\n14 puntos, 0.7% info restante\nAccuracy: 19.5%",
            "V2: Sparsity escalonada\n14 puntos, 9% info restante\nAccuracy: 24.2%",
            "V3: Per-layer tracking\n14 puntos con métricas\nAccuracy: 28.9%",
            "V4: 7 puntos estratégicos\n52% info preservada\nAccuracy: En entrenamiento"
        ]
    },
    22: {  # Conclusiones
        "text_blocks": [
            "HSAQ permite entrenar modelos\n11.5× más grandes sin aumentar\nVRAM de optimizer",
            "Sparsity adaptativa elimina\nnecesidad de AdamW\n(4 bytes/param vs 8)",
            "Arquitectura híbrida funcional:\nTransformer + SNN + SSM + JEPA\ncon HSAQ en cada componente"
        ]
    },
}

def replace_text_in_shape(shape, old_text, new_text):
    """Replace text in a shape, preserving formatting where possible."""
    if not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
        # Also check full paragraph text (for multi-run text boxes)
        if old_text in para.text:
            full = para.text
            for run in para.runs:
                run.text = ""
            # Put all replaced text in first run
            if para.runs:
                para.runs[0].text = full.replace(old_text, new_text)
                replaced = True
    return replaced

def fill_slide_details(slide, details):
    """Fill detailed content into a slide."""
    if "items" in details:
        # For agenda slides with numbered items
        for s in slide.shapes:
            if s.has_text_frame and "04" in s.text_frame.text:
                pass  # Keep slide number
            elif s.has_text_frame and len(s.text_frame.text.strip()) > 15:
                # Try to find text blocks
                pass
        # We'll append items to the main text body
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                t = s.text_frame.text.strip()
                if len(t) > 20 and "AGENDA" not in t and "04" not in t:
                    s.text_frame.paragraphs[0].clear()
                    s.text_frame.paragraphs[0].text = "\n".join(f"   {i+1}. {item}" for i, item in enumerate(details["items"]))
                    break

    if "text_blocks" in details:
        blocks = details["text_blocks"]
        bidx = 0
        for s in slide.shapes:
            if s.has_text_frame:
                t = s.text_frame.text.strip()
                if len(t) > 15 and bidx < len(blocks):
                    s.text_frame.paragraphs[0].clear()
                    s.text_frame.paragraphs[0].text = blocks[bidx]
                    bidx += 1

    if "quote" in details:
        for s in slide.shapes:
            if s.has_text_frame and "talento" in s.text_frame.text.lower():
                s.text_frame.paragraphs[0].clear()
                s.text_frame.paragraphs[0].text = details["quote"]


def generate():
    prs = Presentation(TEMPLATE)
    n_slides = len(prs.slides)
    print(f"Template: {n_slides} slides")

    # ── First pass: replace text markers ──
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for old, new in CONTENT.items():
                replace_text_in_shape(shape, old, new)

    # ── Second pass: fill detailed content ──
    for si, details in CONTENT_DETAILS.items():
        if si <= n_slides:
            fill_slide_details(prs.slides[si - 1], details)

    # ── Remove guide slides (25-33) ──
    # We'll keep slides 1-24 (presentation) and remove the guide section
    # python-pptx doesn't support deleting slides easily
    # Instead, we'll use the XML part relationship approach

    slide_ids_to_remove = list(range(24, n_slides))  # 0-indexed
    if slide_ids_to_remove:
        # Remove from last to first to maintain indices
        for idx in reversed(slide_ids_to_remove):
            if idx < len(prs.slides):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

    prs.save(OUTPUT)
    print(f"Guardado: {OUTPUT}")
    print(f"Slides finales: {len(prs.slides)}")

if __name__ == "__main__":
    generate()
