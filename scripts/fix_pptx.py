"""
Segunda pasada: reemplazar textos pendientes en la presentación PPTX
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

PPTX = "/home/methodwhite/Documentos/IA LATAM/M.A.T.E.R.I.A./MATERIA_HSAQ_Conferencia_AI_LATAM.pptx"

prs = Presentation(PPTX)

def replace_text_in_shape(shape, old_text, new_text):
    if not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
        if old_text in para.text:
            full = para.text
            runs_text = [r.text for r in para.runs]
            for r in para.runs:
                r.text = ""
            if para.runs:
                para.runs[0].text = full.replace(old_text, new_text)
                replaced = True
    return replaced

def set_text(shape, new_text):
    if shape.has_text_frame:
        shape.text_frame.paragraphs[0].clear()
        shape.text_frame.paragraphs[0].text = new_text

def clear_shape(shape):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para.clear()

# ── Slide 2: Title of talk ──
sl2 = prs.slides[1]
for s in sl2.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "título" in t.lower():
            set_text(s, "M.A.T.E.R.I.A. V4\n+ HSAQ")
        if "Nombre" in t:
            set_text(s, "Method White   ·   Investigador · M.A.T.E.R.I.A. Research")

# ── Slides 8-9: Content with image ──
sl8 = prs.slides[7]
for s in sl8.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "Texto a la izquierda" in t:
            set_text(s, "HSAQ reduce la huella de memoria del optimizer en un 50%\n\n"
                        "• AdamW: 2 estados por parámetro (8 bytes/param)\n"
                        "• SGD Nesterov: 1 estado (4 bytes/param)\n"
                        "• En 190M parámetros: ~760MB ahorrados\n"
                        "• Permite modelos 11.5× más grandes en mismo hardware")

sl9 = prs.slides[8]
for s in sl9.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "Imagen a la izquierda" in t:
            set_text(s, "HSAQ no es INT8 ni INT4\n\n"
                        "HSAQ es CUANTIZACIÓN DE ACTIVACIONES mediante sparsity adaptativa.\n\n"
                        "• Máscara binaria {0, valor_original}\n"
                        "• kthvalue por batch → umbral dinámico\n"
                        "• sparsity escalonada por capa (5%→15%)\n"
                        "• 52% de información preservada")

# ── Slide 15: Cita (quote) ──
sl15 = prs.slides[14]
for s in sl15.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "reemplaza" in t.lower() or "talento" in t.lower():
            set_text(s, "La cuantización adaptativa no reemplaza la arquitectura,\n"
                        "la potencia. HSAQ democratiza el acceso a modelos\n"
                        "grandes al reducir los requisitos de hardware.")

# ── Slide 16: Tabla comparativa ──
sl16 = prs.slides[15]
for s in sl16.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "Tabla comparativa" in t:
            set_text(s, "Métrica\tSin HSAQ (AdamW)\tCon HSAQ (SGD)\n"
                        "Params máximos\t16.5M\t190M (11.5×)\n"
                        "Optimizer VRAM\t8 bytes/param\t4 bytes/param\n"
                        "SNN activo\tNo (spike=0)\tSí (spike=0.04)\n"
                        "Info preservada\t100%\t52% (controlada)\n"
                        "Regularización\tWeight decay\tSparsity adaptativa")

# ── Slide 17: Resultados ──
sl17 = prs.slides[16]
for s in sl17.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "Gráfico" in t or "lectura" in t:
            set_text(s, "Accuracy: 28.9% (+48% vs sparsity uniforme)\n"
                        "Perplexity: 49 (vs 107 con sparsity uniforme)\n"
                        "Loss: 3.35 (vs 3.73 con sparsity uniforme)\n"
                        "Training estable sin OOMs por más de 3200 batches")

# ── Slide 19: Checklist ──
sl19 = prs.slides[18]
for s in sl19.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "Lista de verificación" in t:
            set_text(s, "✓ Sparsity escalonada implementada\n"
                        "✓ Per-layer HSAQ tracking funcional\n"
                        "✓ Bug kthvalue corregido\n"
                        "✓ SSM mantiene HSAQ con sparsity 0.05\n"
                        "✓ 7 puntos estratégicos (vs 14 originales)\n"
                        "✓ 52% información preservada (vs 0.7%)")

# ── Slide 21: Participa ──
sl21 = prs.slides[20]
for s in sl21.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "mayor reto" in t.lower():
            set_text(s, "¿Cómo crees que HSAQ puede\n"
                        "aplicarse en tu organización?")

# ── Slide 24: Únete ──
sl24 = prs.slides[23]
for s in sl24.shapes:
    if s.has_text_frame:
        t = s.text_frame.text
        if "Sigue conectado" in t:
            set_text(s, "Sigue el proyecto en github.com/MethodWhite/MATERIA\n"
                        "Documentación completa y paper científico disponibles")

prs.save(PPTX)
print("✅ Textos pendientes actualizados")
