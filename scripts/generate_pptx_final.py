#!/usr/bin/env python3
"""
Genera PPTX presentación MATERIA V4 + HSAQ para Congreso AI LATAM
Usa la plantilla oficial y reemplaza contenido preservando el diseño.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import lxml.etree as etree

TEMPLATE = "/home/methodwhite/Documentos/IA LATAM/Plantilla/Plantilla_Congreso_AI_LATAM.pptx"
OUTPUT   = "/home/methodwhite/Documentos/IA LATAM/M.A.T.E.R.I.A./MATERIA_HSAQ_Conferencia_AI_LATAM.pptx"

def replace_shape_text(shape, old_text, new_text):
    """Replace text in a shape's text frame, preserving formatting."""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    replaced = False
    for para in tf.paragraphs:
        full = para.text
        if old_text in full:
            # Clear all runs and set text on first run
            for run in para.runs:
                run.text = ""
            if para.runs:
                para.runs[0].text = full.replace(old_text, new_text)
            else:
                # Create a run if none exist
                run = para.add_run()
                run.text = full.replace(old_text, new_text)
            replaced = True
    return replaced

def set_shape_text(shape, new_text):
    """Set entire text of a shape's first paragraph."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs:
        # Clear all runs in first paragraph
        p = tf.paragraphs[0]
        for run in p.runs:
            run.text = ""
        if p.runs:
            p.runs[0].text = new_text
        else:
            run = p.add_run()
            run.text = new_text

def remove_slide(prs, slide_index):
    """Remove a slide by index (0-based)."""
    rId = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[slide_index])


def main():
    prs = Presentation(TEMPLATE)
    n = len(prs.slides)
    print(f"Template: {n} slides")

    # ─── CONTENT MAPPING: (template slide index 0-based → {old_text: new_text}) ───
    # Shows what text to replace on each slide

    replacements = {
        # Slide 1 - Portada
        0: {
            "CONGRESO · INTELIGENCIA ARTIFICIAL · 2026":
                "M.A.T.E.R.I.A. V4 · ARQUITECTURA HSAQ · 2026",
            "Congreso": "M.A.T.E.R.I.A.",
            "AI LATAM": "V4 + HSAQ",
            "El encuentro de la comunidad de Inteligencia Artificial de Latinoamérica.":
                "Arquitectura de Cuantización Adaptativa Hiperdispersa para modelos de próxima generación.",
            "Fecha del evento": "30 Julio, 2026",
            "Ciudad · Sede": "Online · LATAM",
            "ialatam.com": "github.com/MethodWhite/MATERIA",
        },

        # Slide 2 - Charla magistral
        1: {
            "CHARLA MAGISTRAL · TRACK DE IA":
                "M.A.T.E.R.I.A. V4 · ARQUITECTURA HSAQ",
            "Título de tu charla en\nuna o dos líneas":
                "M.A.T.E.R.I.A. V4 + HSAQ:\nCuantización Adaptativa",
            "Tu Nombre Apellido   ·   Cargo · Organización":
                "MethodWhite   ·   Investigador · M.A.T.E.R.I.A. Research",
        },

        # Slide 3 - Sección: Introducción
        2: {
            "SECCIÓN": "INTRODUCCIÓN",
            "Título de la sección":
                "M.A.T.E.R.I.A. V4 y el problema de la cuantización",
        },

        # Slide 4 - Agenda
        3: {
            "Lo que veremos hoy": "Agenda",
        },

        # Slide 5 - Sobre el ponente / MATERIA
        4: {
            "SOBRE EL PONENTE": "M.A.T.E.R.I.A. V4",
            "Quién te acompaña hoy": "¿Qué es M.A.T.E.R.I.A.?",
        },

        # Slide 6 - Contenido: Qué es HSAQ
        5: {
            "Título de la lámina de contenido": "¿Qué es HSAQ?",
        },

        # Slide 7 - Comparativa: HSAQ vs TurboQuant
        6: {
            "COMPARATIVA": "HSAQ vs TURBOQUANT",
            "Dos columnas de contenido": "Cuantización adaptativa vs cuantización fija",
        },

        # Slide 8 - Contenido: texto+imagen
        7: {
            "Texto a la izquierda, imagen a la derecha":
                "HSAQ reduce memoria del optimizer en 50%",
        },

        # Slide 9 - Contenido: imagen+texto
        8: {
            "Imagen a la izquierda, texto a la derecha":
                "HSAQ no es INT8 ni INT4",
        },

        # Slide 10 - Conceptos: tres pilares
        9: {
            "CONCEPTOS": "PILARES DE HSAQ",
            "Tres ideas en tarjetas": "Principios fundamentales",
        },

        # Slide 11 - Pilares: aplicación por capas
        10: {
            "PILARES": "APLICACIÓN POR CAPAS",
            "Cuatro bloques en cuadrícula": "Puntos de aplicación HSAQ en el pipeline",
        },

        # Slide 12 - Proceso: pipeline
        11: {
            "PROCESO": "PIPELINE HSAQ",
            "Un proceso en cuatro pasos": "Forward pass con sparsity escalonada",
        },

        # Slide 13 - Comparación: sin HSAQ vs con HSAQ
        12: {
            "COMPARACIÓN": "SIN HSAQ VS CON HSAQ",
            "Enfoque tradicional vs. con IA": "AdamW vs SGD Nesterov + HSAQ",
        },

        # Slide 14 - En cifras: métricas
        13: {
            "EN CIFRAS": "MÉTRICAS Y RESULTADOS",
            "Los datos que respaldan tu mensaje": "190M parámetros entrenando en RTX 3050",
        },

        # Slide 15 - Cita
        14: {
            "Autor de la cita  ·  Cargo / referencia":
                "MethodWhite  ·  M.A.T.E.R.I.A. Research 2026",
        },

        # Slide 16 - Datos: tabla
        15: {
            "DATOS": "COMPARATIVA",
            "Tabla comparativa": "Sin HSAQ (AdamW) vs Con HSAQ (SGD)",
        },

        # Slide 17 - Resultados: gráfico
        16: {
            "RESULTADOS": "RESULTADOS",
            "Gráfico con su lectura clave": "Accuracy: 28.9% · Perplexity: 31.7 · +48% mejora",
        },

        # Slide 18 - Arquitectura
        17: {
            "ARQUITECTURA": "ARQUITECTURA",
            "Diagrama de flujo o arquitectura":
                "Input → Emb → HSAQ → Transformer ×10 → SNN → HSAQ → SSM → HSAQ → JEPA → HSAQ → Head",
        },

        # Slide 19 - Checklist
        18: {
            "CHECKLIST": "LOGROS HSAQ",
            "Lista de verificación": "7 puntos estratégicos · Sparsity escalonada · Per-layer tracking",
        },

        # Slide 20 - Trayectoria
        19: {
            "TRAYECTORIA": "EVOLUCIÓN HSAQ",
            "Línea de tiempo del proyecto": "V1 → V2 → V3 → V4: 0.7% → 52% info preservada",
        },

        # Slide 21 - Participa
        20: {
            "PARTICIPA": "PARTICIPA",
        },

        # Slide 22 - Para llevar
        21: {
            "PARA LLEVAR": "CONCLUSIONES",
            "Tres ideas para recordar": "Ideas clave",
        },

        # Slide 23 - Gracias
        22: {
            "¿Conversamos? Espacio para preguntas.":
                "¿Conversamos? methodwhite@github.com",
            "tu@correo.com      ·      @tu_usuario":
                "methodwhite@github.com  ·  @MethodWhite",
        },

        # Slide 24 - Únete
        23: {
            "ÚNETE A LA COMUNIDAD": "CONTACTO",
            "Sigue conectado con":
                "Sigue el proyecto en github.com/MethodWhite/MATERIA",
        },
    }

    # ─── DETAILED CONTENT (body text) ───
    # These are the main text blocks for each content slide

    detailed_content = {
        # Slide 4 - Agenda: add bullet items to the main text box
        3: [
            "1.  ¿Qué es M.A.T.E.R.I.A. V4?\n"
            "2.  El problema de la cuantización\n"
            "3.  HSAQ: cuantización adaptativa\n"
            "4.  Arquitectura por capas\n"
            "5.  Resultados y métricas\n"
            "6.  Próximos pasos",
        ],

        # Slide 5 - Qué es MATERIA V4
        4: [
            "M.A.T.E.R.I.A. es un modelo híbrido que integra cuatro paradigmas de "
            "procesamiento neuronal:\n\n"
            "• Transformers con Flash Attention 2 + RoPE + GQA\n"
            "• LIF-SNN (neuronas Leaky Integrate-and-Fire)\n"
            "• State Space Models (SSM)\n"
            "• JEPA (Joint Embedding Predictive Architecture)\n\n"
            "Entrenado con HSAQ: cuantización adaptativa que reemplaza a AdamW "
            "como optimizador, reduciendo la memoria en un 50%.",
        ],

        # Slide 6 - Qué es HSAQ
        5: [
            "HyperSparse Adaptive Quantization\n\n"
            "HSAQ no es INT8 ni INT4. HSAQ es cuantización de activaciones "
            "mediante sparsity adaptativa dinámica.\n\n"
            "• kthvalue → umbral dinámico por batch\n"
            "• Máscara binaria {0, valor_original}\n"
            "• Sin calibración, sin datasets externos\n"
            "• Hardware-agnostic (CPU/GPU/TPU)\n"
            "• 7 puntos estratégicos con sparsity escalonada",
        ],

        # Slide 7 - HSAQ vs TurboQuant (left column / main)
        6: [
            "TurboQuant (Google):\n\n"
            "• Cuantización fija INT8 de pesos\n"
            "• Requiere calibración post-entrenamiento\n"
            "• Solo GPU con soporte INT8\n"
            "• No adaptativo\n\n"
            "HSAQ (MATERIA):\n\n"
            "• Sparsity adaptativa por batch\n"
            "• Sin calibración\n"
            "• CPU/GPU/TPU sin modificaciones\n"
            "• Umbral dinámico por capa",
        ],

        # Slide 8 - Memoria optimizer
        7: [
            "HSAQ reduce la huella de memoria del optimizer en un 50%\n\n"
            "• AdamW: 2 estados FP32 por parámetro (8 bytes/param)\n"
            "• SGD Nesterov: 1 estado FP32 (4 bytes/param)\n"
            "• En 190M parámetros: ~760 MB ahorrados\n"
            "• Permite modelos 11.5× más grandes en mismo hardware\n"
            "• Todo en RTX 3050 (4 GB VRAM)",
        ],

        # Slide 9 - No es INT8/INT4
        8: [
            "HSAQ es CUANTIZACIÓN DE ACTIVACIONES\n\n"
            "• Máscara binaria: {0, valor_original}\n"
            "• kthvalue por batch → umbral dinámico\n"
            "• Sparsity escalonada por capa (5% → 15%)\n"
            "• 52% de información preservada\n"
            "• Sin INT8, sin INT4, sin weight quantization\n\n"
            "HSAQ optimiza el uso de recursos haciendo que "
            "solo las neuronas relevantes participen en el cómputo.",
        ],

        # Slide 10 - Tres pilares
        9: [
            "ADAPTATIVO\nUmbral recalculado\nen cada batch vía\nkthvalue dinámico",
            "POR CAPAS\nCada componente\ntiene su propia\nsparsity calibrada",
            "SIN ESTADO\nNo hay buffers\npersistentes entre\nbatches de training",
        ],

        # Slide 11 - Aplicación por capas
        10: [
            "EMBEDDING\nSparsity: 5%\nPreserva 95% de\ninformación de entrada",
            "TRANSFORMER (c/3)\nSparsity: 8-15%\nProgresivo según\nprofundidad",
            "SNN + SSM\nSparsity: 10% y 5%\nProcesamiento\ntemporal controlado",
            "JEPA LATENT\nSparsity: 5%\nPreserva espacio\nde predicción",
        ],

        # Slide 12 - Pipeline
        11: [
            "1. FORWARD PASS\nEmbedding → HSAQ(5%) →\nTransformer → HSAQ(8-15%) →\nSNN → HSAQ(10%) → SSM → HSAQ(5%) →\nJEPA → HSAQ(5%) → Head",
            "2. SPARSITY ESCALONADA\n0.05 → 0.15 progresivo\n7 puntos estratégicos\n(antes 14, ahora óptimo)",
            "3. BACKWARD PASS\nGradiente fluye solo por\nneuronas activas\n(STE: Straight-Through Estimator)",
            "4. WEIGHT UPDATE\nSGD Nesterov\nmomentum = 0.9\nlr = 5×10⁻⁴",
        ],

        # Slide 13 - Sin HSAQ vs Con HSAQ
        12: [
            "SIN HSAQ (AdamW)\n\n"
            "• 100% neuronas activas\n"
            "• Mayor sobreajuste\n"
            "• 8 bytes/param optimizer\n"
            "• SNN desactivado (spike=0)\n"
            "• Máx: 16.5M parámetros\n"
            "• Sin tracking por capa",
            "CON HSAQ (SGD)\n\n"
            "• 70% activas (regulariza)\n"
            "• Menos overfitting\n"
            "• 4 bytes/param (-50% VRAM)\n"
            "• SNN activo (spike=0.04)\n"
            "• 190M parámetros (11.5×)\n"
            "• Per-layer tracking ✓",
        ],

        # Slide 14 - Métricas
        13: [
            "190M\nparámetros",
            "11.5×\nmás params",
            "34%\naccuracy",
            "+48%\nmejora",
            "52%\ninfo preservada",
            "31.7\nperplexity",
        ],

        # Slide 15 - Cita (replace main quote)
        14: [
            "La cuantización adaptativa no reemplaza la arquitectura, "
            "la potencia. HSAQ democratiza el acceso a modelos grandes "
            "al reducir los requisitos de hardware a la mitad.",
        ],

        # Slide 16 - Tabla comparativa
        15: [
            "Métrica\tSin HSAQ (AdamW)\tCon HSAQ (SGD)\n"
            "Params máx.\t16.5M\t190M (11.5×)\n"
            "Optimizer RAM\t8 bytes/param\t4 bytes/param\n"
            "SNN activo\tNo (spike=0)\tSí (spike>0)\n"
            "Info preservada\t100%\t52% (controlada)\n"
            "Tracking\tNo\tPer-layer HSAQ ✓",
        ],

        # Slide 17 - Resultados
        16: [
            "Accuracy: 28.9% (+48% vs sparsity uniforme 19.5%)\n"
            "Perplexity: 31.7 (vs 107 con sparsity uniforme)\n"
            "Loss: 3.30 (vs 3.73 con sparsity uniforme)\n"
            "3200+ batches sin OOMs\n"
            "Per-layer sparsity targets exactos en todos los puntos",
        ],

        # Slide 18 - Arquitectura (detailed)
        17: [
            "Input Token Embedding\n→\nHSAQ (sparsity=5%)\n→\nBlocks Transformer ×10\n"
            "c/ Flash Attn 2 + RoPE + GQA\n→\nHSAQ (sparsity=8-15% c/3)\n"
            "→\nLIF-SNN → HSAQ(10%)\n→\nSSM → HSAQ(5%)\n→\nJEPA → HSAQ(5%)\n→\nLayerNorm → Head",
        ],

        # Slide 19 - Checklist
        18: [
            "✓ Sparsity escalonada implementada (5%→15% progresivo)\n"
            "✓ Per-layer HSAQ tracking funcional\n"
            "✓ Bug kthvalue corregido (k = n × sparsity)\n"
            "✓ SSM mantiene HSAQ con sparsity 0.05\n"
            "✓ 7 puntos estratégicos (vs 14 originales)\n"
            "✓ 52% información preservada (vs 0.7%)",
        ],

        # Slide 20 - Evolución
        19: [
            "V1 · Sparsity uniforme\n14 pts, 0.7% info\nAccuracy: 19.5%",
            "V2 · Sparsity escalonada\n14 pts, 9% info\nAccuracy: 24.2%",
            "V3 · Per-layer tracking\n14 pts con métricas\nAccuracy: 28.9%",
            "V4 · 7 pts estratégicos\n52% info preservada\nAccuracy: 34%+",
        ],

        # Slide 21 - Participa
        20: [
            "¿Cómo crees que la cuantización adaptativa\n"
            "puede aplicarse en tu organización?",
        ],

        # Slide 22 - Conclusiones
        21: [
            "HSAQ permite entrenar modelos 11.5× más grandes\n"
            "sin aumentar VRAM de optimizer",
            "Sparsity adaptativa elimina necesidad de AdamW:\n"
            "4 bytes/param vs 8 bytes/param",
            "Arquitectura híbrida funcional:\n"
            "Transformer + SNN + SSM + JEPA",
        ],

        # Slide 24 - Contacto
        23: [
            "Documentación completa y paper científico\n"
            "disponibles en github.com/MethodWhite/MATERIA\n\n"
            "Código abierto · Comunidad · Investigación",
        ],
    }

    # ─── APPLY REPLACEMENTS ───
    for slide_idx, replaces in replacements.items():
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for old_text, new_text in replaces.items():
                replace_shape_text(shape, old_text, new_text)

    # ─── APPLY DETAILED CONTENT ───
    for slide_idx, text_blocks in detailed_content.items():
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        block_idx = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            # Skip small shapes (decorators, slide numbers, logos)
            if len(t) < 5:
                continue
            # Skip header/title texts that were already replaced
            skip_patterns = [
                "M.A.T.E.R.I.A.", "HSAQ", "CONTENIDO", "COMPARATIVA", "CONCEPTOS",
                "PILARES", "PROCESO", "COMPARACIÓN", "EN CIFRAS", "DATOS",
                "RESULTADOS", "ARQUITECTURA", "CHECKLIST", "TRAYECTORIA",
                "PARTICIPA", "PARA LLEVAR", "CONCLUSIONES", "INTRODUCCIÓN",
                "AGENDA", "SOBRE EL PONENTE", "LOGROS", "EVOLUCIÓN",
                "CONTACTO", "¡Gracias!", "ÚNETE", "APLICACIÓN",
                "PILARES DE", "SIN HSAQ", "CON HSAQ", "MÉTRICAS",
                "COMPARATIVA", "IDEAS CLAVE", "Agenda",
                "¿Conversamos?", "methodwhite", "@MethodWhite",
                "01", "02", "03", "04", "05", "06", "07", "08", "09",
                "10", "11", "12", "13", "14", "15", "16", "17",
                "18", "19", "20", "21", "22", "23", "24",
                "30 Julio", "Online", "github.com", "IA LATAM", "Congreso",
                "MethodWhite", "Investigador", "Arquitectura de Cuantización",
            ]
            is_header = any(t.startswith(sp) for sp in skip_patterns) or len(t) > 200
            if is_header or t in skip_patterns:
                continue
            if block_idx < len(text_blocks):
                set_shape_text(shape, text_blocks[block_idx])
                block_idx += 1

    # ─── REMOVE GUIDE SLIDES (25-33, 0-indexed: 24-32) ───
    for idx in range(32, 23, -1):  # Remove from last to first
        if idx < len(prs.slides):
            remove_slide(prs, idx)
            print(f"  Removed slide {idx+1}")

    # ─── SAVE ───
    prs.save(OUTPUT)
    print(f"\n✅ PPTX guardado: {OUTPUT}")
    print(f"   Slides finales: {len(prs.slides)}")


if __name__ == "__main__":
    main()
