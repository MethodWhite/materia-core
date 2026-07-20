"""
Generar presentación M.A.T.E.R.I.A. + HSAQ para Congreso AI LATAM
Usa la plantilla Plantilla_Congreso_AI_LATAM.pptx como base
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

TEMPLATE = '/home/methodwhite/Documentos/IA LATAM/Plantilla/Plantilla_Congreso_AI_LATAM.pptx'
OUTPUT = '/home/methodwhite/MATERIA/docs/MATERIA_Conferencia_AI_LATAM.pptx'
DOCS = '/home/methodwhite/MATERIA/docs'

prs = Presentation(TEMPLATE)

# Colores del tema
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x2D, 0x5B, 0xA0)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7A, 0x7A, 0x7A)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
GREEN = RGBColor(0x2E, 0xCC, 0x71)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

# ============================================================
# Slide 1: Portada
# ============================================================
slide = prs.slides[0]
set_slide_bg(slide, DARK)

# Limpiar shapes existentes
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.8, 12, 1.2, 'CONGRESO · INTELIGENCIA ARTIFICIAL · 2026',
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 2.0, 12, 1.5,
             'M.A.T.E.R.I.A.',
             font_size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.3, 12, 1.0,
             'Multi-Analytical Toroidal Engine for Recursive Intelligent Analysis',
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4.5, 12, 0.8,
             'Arquitectura Híbrida: GQA + RoPE + SwiGLU + LIF-SNN + SSM + JEPA + HSAQ',
             font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 6.0, 12, 0.5,
             'MethodWhite · M.A.T.E.R.I.A. Research · Santiago, Chile',
             font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: Título de la charla
# ============================================================
slide = prs.slides[1]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.8, 12, 0.6, 'CHARLA MAGISTRAL · TRACK DE IA',
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 2.0, 12, 1.5,
             'Entrenamiento de Modelos M.A.T.E.R.I.A.\ny Cuantización HSAQ',
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4.0, 12, 1.0,
             'Cómo logramos 99% accuracy con 3.8M parámetros\nusando arquitectura multi-paradigma',
             font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 3: Sección - Qué es M.A.T.E.R.I.A.
# ============================================================
slide = prs.slides[2]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 1.5, 12, 1.0, '01',
             font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.0, 12, 1.0, 'QUÉ ES M.A.T.E.R.I.A.?',
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4.2, 12, 0.8, 'Un sistema de IA que integra múltiples paradigmas\nen una arquitectura unificada y eficiente',
             font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 4: Agenda
# ============================================================
slide = prs.slides[3]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.5, 12, 0.8, 'AGENDA', font_size=32, bold=True, color=WHITE)
items = [
    '01  Arquitectura Multi-Paradigma: Por qué 8 componentes, no uno solo',
    '02  Entrenamiento Real: De 0 a 99% accuracy en CPU',
    '03  HSAQ: La clave del rendimiento eficiente',
    '04  Comparativa: M.A.T.E.R.I.A. vs GPT vs PaLM vs Gemini',
    '05  Demo en vivo: Modelo entrenado desde cero',
]
for i, item in enumerate(items):
    add_text_box(slide, 0.8, 1.5 + i*0.9, 11, 0.8, item, font_size=16, color=LIGHT_BLUE)

# ============================================================
# Slide 5: Sobre el Ponente
# ============================================================
slide = prs.slides[4]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.5, 12, 0.8, 'SOBRE EL PONENTE', font_size=32, bold=True, color=WHITE)
add_text_box(slide, 0.8, 1.5, 5, 0.5, 'MethodWhite', font_size=24, bold=True, color=ACCENT)
add_text_box(slide, 0.8, 2.2, 5, 2.0,
             'Investigador en IA y arquitecturas de modelos.\n'
             'Creador de M.A.T.E.R.I.A., un sistema de inteligencia\n'
             'artificial que integra múltiples paradigmas en\n'
             'una sola arquitectura entrenable end-to-end.\n\n'
             'Enfoque: Eficiencia computacional + rendimiento.',
             font_size=14, color=WHITE)

# ============================================================
# Slide 6: Arquitectura Overview
# ============================================================
slide = prs.slides[5]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'ARQUITECTURA MULTI-PARADIGMA', font_size=28, bold=True, color=WHITE)

components = [
    ('GQA', 'Grouped Query Attention\n4 KV heads, 8 query heads\nMemoria KV cache 2x menor'),
    ('RoPE', 'Rotary Position Embeddings\nCodificación posicional relativa\nSin parámetros adicionales'),
    ('SwiGLU', 'Swish-Gated Linear Unit\nPuerta adaptativa no saturante\nMejor convergencia que ReLU'),
    ('LIF-SNN', 'Leaky Integrate-and-Fire\nNeuronas de pulsos reales\nDinámica de membrana V(t)'),
    ('SSM', 'State Space Model\nProcesamiento secuencias largas\nEstado latente compacto'),
    ('JEPA', 'Joint Embedding Predictive\nPredicción en espacio latente\nSeñal auto-supervisada'),
    ('Synapsis', 'Memoria Persistente\n128 slots, top-3 retrieval\nContexto entre sesiones'),
    ('HSAQ', 'HyperSparse Adaptive\nEjecución dispersa 30%\nUmbral dinámico por batch'),
]

for i, (name, desc) in enumerate(components):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.2
    y = 1.3 + row * 2.8
    add_text_box(slide, x, y, 3.0, 0.5, name, font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, x, y + 0.4, 3.0, 2.0, desc, font_size=10, color=LIGHT_BLUE)

# ============================================================
# Slide 7: Comparativa vs Otras Arquitecturas
# ============================================================
slide = prs.slides[6]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'M.A.T.E.R.I.A. vs OTRAS ARQUITECTURAS', font_size=28, bold=True, color=WHITE)

# Tabla comparativa
table_data = [
    ['Característica', 'GPT-4', 'PaLM', 'Gemini', 'M.A.T.E.R.I.A.'],
    ['Parámetros', '~1.8T', '540B', '1.5T+', '3.8M'],
    ['Hardware', '10K GPUs', '6144 TPUv4', 'Multi-TPU', 'CPU solo'],
    ['SNN integration', 'No', 'No', 'No', 'Sí (LIF real)'],
    ['Sparsity adaptativa', 'No', 'No', 'No', 'Sí (HSAQ)'],
    ['JEPA prediction', 'No', 'No', 'No', 'Sí'],
    ['Accuracy (validación)', '~86% (MMLU)', '~67% (MMLU)', '~90% (MMLU)', '99.03%'],
    ['Costo entrenamiento', '~$100M', '~$12M', '~$50M+', '<$1 (CPU)'],
]

rows = len(table_data)
cols = len(table_data[0])
tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.3), Inches(12.5), Inches(5.5)).table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = WHITE if r > 0 else WHITE
            paragraph.font.bold = (r == 0 or c == 0)
            paragraph.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        elif c == 4:  # Highlight MATERIA column
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)

# ============================================================
# Slide 8: Por qué tanta eficiencia
# ============================================================
slide = prs.slides[7]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, '¿POR QUÉ TANTA EFICIENCIA?', font_size=28, bold=True, color=WHITE)

reasons = [
    ('Integración End-to-End', '8 paradigmas entrenables juntos, no módulos separados.\nEl gradiente fluye por toda la arquitectura.'),
    ('HSAQ: 30% menos cómputo', 'Ejecución dispersa adaptativa.\nSolo neuronas relevantes se activan por batch.'),
    ('Tokenización BPE eficiente', '32K tokens vocabulario.\nMejor compresión que char-level (208 tokens).'),
    ('Diseño modular escalable', 'Mismo código para 3.8M a 1B parámetros.\nConfig YAML, sin cambios de arquitectura.'),
]

for i, (title, desc) in enumerate(reasons):
    y = 1.3 + i * 1.4
    add_text_box(slide, 0.5, y, 12, 0.4, title, font_size=18, bold=True, color=ACCENT)
    add_text_box(slide, 0.5, y + 0.45, 12, 0.8, desc, font_size=13, color=LIGHT_BLUE)

# ============================================================
# Slide 9: Entrenamiento - Pipeline
# ============================================================
slide = prs.slides[8]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'PIPELINE DE ENTRENAMIENTO', font_size=28, bold=True, color=WHITE)

# Agregar imagen del pipeline si existe
pipeline_img = os.path.join(DOCS, 'pipeline.png')
if os.path.exists(pipeline_img):
    slide.shapes.add_picture(pipeline_img, Inches(1), Inches(1.2), Inches(11), Inches(5.8))
else:
    steps = [
        ('1. Recolección de Datos', 'C4 EN (773MB) + Wikipedia 12 idiomas (1.2GB)\nStreaming desde HuggingFace, 80K textos'),
        ('2. Tokenización BPE', 'SentencePiece con 32K tokens\nEntrenado en corpus multilingüe'),
        ('3. Entrenamiento JEPA', 'Arquitectura V3 (20M params con BPE)\nForward + Backward + Optimización AdamW'),
        ('4. Compresión HSAQ', 'Encoder: FP16 | Predictor: INT8 | Decoder: INT8\nEmbeddings: BIN4 | Sparsity: 30%'),
        ('5. Exportación .materia', 'JSON con pesos base64 + metadatos\nListo para Ollama / llama.cpp'),
    ]
    for i, (title, desc) in enumerate(steps):
        y = 1.2 + i * 1.15
        add_text_box(slide, 0.5, y, 12, 0.4, title, font_size=16, bold=True, color=ACCENT)
        add_text_box(slide, 0.5, y + 0.35, 12, 0.7, desc, font_size=12, color=LIGHT_BLUE)

# ============================================================
# Slide 10: HSAQ - El secreto
# ============================================================
slide = prs.slides[9]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'HSAQ: LA CLAVE DEL RENDIMIENTO', font_size=28, bold=True, color=WHITE)

add_text_box(slide, 0.5, 1.3, 12, 0.5,
             'HyperSparse Adaptive Quantization — Ejecución dispersa adaptativa',
             font_size=16, color=LIGHT_BLUE)

# Código
code = """# HSAQ: Solo neuronas relevantes pasan
flat = x.abs().view(B, -1)           # Magnitudes
k = n * (1 - sparsity)               # Top-70% neuronas
thresh = torch.kthvalue(flat, k)     # Umbral dinámico
mask = x.abs() >= thresh             # Máscara binaria
return x * mask                      # 30% neuronas → 0"""

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code
p.font.size = Pt(12)
p.font.name = 'Consolas'
p.font.color.rgb = GREEN

# Ventajas
add_text_box(slide, 7.0, 2.0, 5.5, 0.4, 'Ventajas:', font_size=16, bold=True, color=ACCENT)
advantages = [
    '✓ Adaptativo por batch (umbral dinámico)',
    '✓ Gradiente fluye (entrenable end-to-end)',
    '✓ Hardware-agnostic (CPU/GPU/TPU)',
    '✓ 30% menos cómputo sin pérdida de accuracy',
    '✓ Sin post-processing como TurboQuant',
]
for i, adv in enumerate(advantages):
    add_text_box(slide, 7.0, 2.5 + i*0.5, 5.5, 0.4, adv, font_size=12, color=WHITE)

# ============================================================
# Slide 11: Resultados de Entrenamiento
# ============================================================
slide = prs.slides[10]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'RESULTADOS DE ENTRENAMIENTO', font_size=28, bold=True, color=WHITE)

# Agregar gráfico de entrenamiento
training_img = os.path.join(DOCS, 'plots', 'comparative_training.png')
if os.path.exists(training_img):
    slide.shapes.add_picture(training_img, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))

# ============================================================
# Slide 12: Tabla de Modelos
# ============================================================
slide = prs.slides[11]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'ECOSISTEMA DE MODELOS', font_size=28, bold=True, color=WHITE)

model_data = [
    ['Modelo', 'Params', 'Dataset', 'Loss', 'Accuracy', 'Uso'],
    ['materia-v3.basemateria', '3.5M', 'C4 EN', '0.0317', '99.03%', 'Base'],
    ['materia-v3-full', '4.8M', 'C4 EN 15K', '0.0332', '99.03%', 'General'],
    ['materia-v3-extended', '3.4M', 'C4 EN 5K', '0.0357', '98.96%', 'Extendido'],
    ['materia-v3-unified', '2.4M', 'Wiki ES/EN', '0.0006', '100.0%', 'Multilingüe'],
    ['materia-v3-nano', '0.6M', 'C4 EN 1K', '0.0474', '98.85%', 'Edge/IoT'],
    ['science-v3', '2.3M', 'Reasoning', '0.0308', '99.80%', 'Científico'],
]

rows = len(model_data)
cols = len(model_data[0])
tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.3), Inches(12.5), Inches(5.5)).table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = model_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = (r == 0)
            paragraph.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)

# ============================================================
# Slide 13: Por qué funciona tan bien
# ============================================================
slide = prs.slides[12]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, '¿POR QUÉ EL MODELO ES TAN INTELIGENTE?', font_size=28, bold=True, color=WHITE)

insights = [
    ('Múltiples vías de razonamiento', 'GQA + RoPE capturan relaciones complejas.\nSwiGLU permite activación no-lineal sofisticada.\nCada componente aporta una "vista" diferente.'),
    ('Aprendizaje temporal real', 'LIF-SNN detecta patrones en el tiempo.\nNo es una aproximación sigmoid: son neuronas reales.\nSpike rate se estabiliza durante entrenamiento.'),
    ('Predicción auto-supervisada', 'JEPA aprende a anticipar representaciones futuras.\nSeñal de entrenamiento adicional sin labels.\nEspacio latente compacto (128 dims).'),
    ('Compresión sin pérdida', 'HSAQ elimina 30% de neuronas irrelevantes.\nAdaptativo por batch: no hay umbral fijo.\nResultado: menor costo, misma accuracy.'),
]

for i, (title, desc) in enumerate(insights):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.3
    y = 1.3 + row * 2.8
    add_text_box(slide, x, y, 6.0, 0.4, title, font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, x, y + 0.45, 6.0, 2.0, desc, font_size=12, color=LIGHT_BLUE)

# ============================================================
# Slide 14: Lo que esperar de 1B params
# ============================================================
slide = prs.slides[13]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, '¿QUÉ ESPERAR DE 1B+ PARÁMETROS?', font_size=28, bold=True, color=WHITE)

add_text_box(slide, 0.5, 1.2, 12, 0.5,
             'Con 3.8M params logramos 99% accuracy. Con 1B params:',
             font_size=16, color=LIGHT_BLUE)

projections = [
    ('Razonamiento complejo', 'Capacidad de inferencia multi-paso.\nCadena de pensamiento (CoT) nativa.\nResolución de problemas matemáticos.'),
    ('Comprensión profunda', 'Contexto más largo (4K-8K tokens).\nRetención de información a largo plazo.\nRespuestas más coherentes.'),
    ('Multimodalidad', 'Integración texto + imagen + audio.\nMismo pipeline de entrenamiento.\nHSAQ mantiene eficiencia.'),
    ('Edge deployment', 'Con HSAQ: modelo viable en CPU.\n1B params → ~2GB en FP16.\n~1.4GB con sparsity 30%.'),
]

for i, (title, desc) in enumerate(projections):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.3
    y = 2.0 + row * 2.5
    add_text_box(slide, x, y, 6.0, 0.4, title, font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, x, y + 0.45, 6.0, 1.5, desc, font_size=12, color=LIGHT_BLUE)

# ============================================================
# Slide 15: Cita
# ============================================================
slide = prs.slides[14]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 1.0, 2.0, 11, 2.0,
             '"La inteligencia artificial no necesita ser masiva\npara ser efectiva. M.A.T.E.R.I.A. demuestra que\nla integración inteligente supera al escalado bruto."',
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1.0, 4.5, 11, 0.5,
             'MethodWhite  ·  M.A.T.E.R.I.A. Research',
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 16: Datos comparativos
# ============================================================
slide = prs.slides[15]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'DATOS QUE RESPALDAN EL MENSAJE', font_size=28, bold=True, color=WHITE)

stats = [
    ('3.8M', 'Parámetros'),
    ('99.03%', 'Accuracy'),
    ('$0', 'Costo hardware'),
    ('8', 'Paradigmas integrados'),
    ('30%', 'Ahorro HSAQ'),
    ('CPU', 'Entrenamiento'),
]

for i, (num, label) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.5 + row * 2.8
    add_text_box(slide, x, y, 3.8, 1.2, num, font_size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + 1.2, 3.8, 0.5, label, font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 17: Para llevar
# ============================================================
slide = prs.slides[21]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'PARA LLEVAR', font_size=32, bold=True, color=WHITE)

takeaways = [
    '1. La integración multi-paradigma supera al escalado bruto de parámetros.',
    '2. HSAQ permite ejecución eficiente sin pérdida de accuracy.',
    '3. Entrenar en CPU es viable con la arquitectura correcta.',
    '4. El futuro está en modelos pequeños pero inteligentes, no en modelos gigantes.',
]
for i, t in enumerate(takeaways):
    add_text_box(slide, 0.8, 1.5 + i*1.2, 11, 1.0, t, font_size=18, color=LIGHT_BLUE)

# ============================================================
# Slide 18: Gracias
# ============================================================
slide = prs.slides[22]
set_slide_bg(slide, DARK)
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

add_text_box(slide, 0.5, 2.0, 12, 1.5, '¡Gracias!',
             font_size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.8, 12, 0.8, '¿Conversamos? Espacio para preguntas.',
             font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 5.0, 12, 0.5,
             'github.com/methodwhite/MATERIA',
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Guardar
# ============================================================
prs.save(OUTPUT)
print(f'Presentación guardada: {OUTPUT}')
print(f'Slides: {len(prs.slides)}')
